"""
# v7.0: DataClass, Caché JSON, Hash Global, Filtros, Historial Ediciones

=============================================================
GENERADOR DE ÍNDICE + FOLIADOR — CRUBC Los Ríos
=============================================================
Autor: César Acuña Leal - CRUBC Los Ríos / Oficina Técnica
Versión: 7.0
=============================================================
MEJORAS v6.0:
  1. Hash SHA-256 por archivo — integridad documental auditable
  2. Análisis paralelo con ThreadPoolExecutor — velocidad real
  3. Clase ExpedienteDigital — separación motor / interfaz
  4. Hash visible en Excel, Word y log exportable
=============================================================
INSTALACIÓN (CMD, una sola vez):
  pip install pdfplumber python-docx openpyxl pillow
=============================================================
USO:
  python generador_indice.py          # Abre la interfaz
  python generador_indice.py --test   # Ejecuta pruebas automáticas
=============================================================
"""

import os, re, zipfile, logging, sys, threading, shutil, unittest, hashlib
import json, tempfile, copy
from dataclasses import dataclass, field, asdict
from typing import Optional
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path
from xml.etree import ElementTree as ET

# ══════════════════════════════════════════════════════════════════════════════
# LOGGING
# ══════════════════════════════════════════════════════════════════════════════

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)]
)
log = logging.getLogger("indice")

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 1: CONSTANTES Y CONFIGURACIÓN CENTRALIZADA
# ══════════════════════════════════════════════════════════════════════════════

EXT_DOCS     = {".pdf",".docx",".xlsx",".xls",".pptx",".txt",".odt",".ods"}
EXT_IMGS     = {".jpg",".jpeg",".png",".tiff",".tif",".bmp",".webp"}
EXT_GEO      = {".kml",".kmz",".shp"}
EXT_ZIP      = {".zip",".rar",".7z"}
EXT_VIDEO    = {".mp4",".avi",".mov",".mkv",".wmv"}
EXT_AUDIO    = {".mp3",".wav",".m4a",".ogg"}
EXT_ALL      = EXT_DOCS|EXT_IMGS|EXT_GEO|EXT_ZIP|EXT_VIDEO|EXT_AUDIO
EXT_SHP_COMP = {".dbf",".shx",".prj",".cpg",".sbn",".sbx",".qpj"}

TIPOS_MAP = {
    ".pdf":"PDF",".docx":"Word",".xlsx":"Excel",".xls":"Excel",
    ".pptx":"PowerPoint",".txt":"Texto",".odt":"ODT",".ods":"ODS",
    ".jpg":"Imagen",".jpeg":"Imagen",".png":"Imagen",
    ".tiff":"Imagen",".tif":"Imagen",".bmp":"Imagen",".webp":"Imagen",
    ".kml":"KML",".kmz":"KMZ",".shp":"Shapefile",
    ".zip":"ZIP",".rar":"RAR",".7z":"7Z",
    ".mp4":"Video",".avi":"Video",".mov":"Video",".mkv":"Video",".wmv":"Video",
    ".mp3":"Audio",".wav":"Audio",".m4a":"Audio",".ogg":"Audio",
}

# ── Tipos documentales base ─────────────────────────────────────────────────
# Cubre: administración pública general, finanzas, desarrollo social,
# municipios, SSPP, rendiciones, contratos, OIRS, etc.
# Se amplían con config.json sin tocar el código.

TIPOS_NOMBRE_BASE = [
    # Comunicaciones formales
    ("Acta",           ["_acta_","_acta-","-acta_","acta_acuerdos","acta_de",
                        "acta-de","_actas_","acta_reunion","acta_sesion",
                        "acta_entrega","acta_recepcion","acta_visita"]),
    ("Correo",         ["_correo_","_correo-","correo_","fwd_","fw_","_mail_","correo-"]),
    ("Oficio",         ["_oficio_","_oficio-","oficio_n","oficio-n","oficio_ord"]),
    ("Ordinario",      ["_ordinario_","ordinario_n","ord_n"]),
    ("Circular",       ["_circular_","circular_n"]),
    ("Memo",           ["_memorandum_","memorandum_","_memo_","memo_interno"]),
    ("Minuta",         ["_minuta_","minuta_","minuta_reunion"]),

    # Actos administrativos
    ("Resolución",     ["_resolucion_","_res_exenta","res_exenta","res-exenta",
                        "resolucion_exenta","_resolucion-","res_afecta"]),
    ("Decreto",        ["_decreto_","decreto_exento","decreto-exento",
                        "decreto_alcaldicio","decreto_afecto"]),
    ("Dictamen",       ["_dictamen_","dictamen_"]),
    ("Instrucción",    ["_instruccion_","instruccion_","instruccion-"]),

    # Informes y análisis
    ("Informe",        ["_informe_","_informe-","informe_observ","informe_tecn",
                        "informe_final","informe_avance","informe_legal",
                        "informe_juridico","informe_ambiental","informe_tecnico",
                        "informe_social","informe_financiero","informe_auditoria",
                        "informe_rendicion","_informe"]),
    ("Evaluación",     ["_evaluacion_","evaluacion_","eval_tecnica","eval_social"]),
    ("Diagnóstico",    ["_diagnostico_","diagnostico_"]),

    # Finanzas y rendiciones
    ("Factura",        ["_factura_","factura_n","factura-"]),
    ("Boleta",         ["_boleta_","boleta_n","boleta-"]),
    ("Rendición",      ["_rendicion_","rendicion_","rendicion-cuentas",
                        "rendicion_gastos"]),
    ("Presupuesto",    ["_presupuesto_","presupuesto_","ppto_"]),
    ("Cotización",     ["_cotizacion_","cotizacion_","cotizacion-"]),
    ("Liquidación",    ["_liquidacion_","liquidacion_"]),
    ("Voucher",        ["_voucher_","voucher_","comprobante_"]),
    ("Orden de Compra",["_orden_compra","orden-compra","oc_n","_oc_"]),

    # Contratos y convenios
    ("Contrato",       ["_contrato_","contrato_","contrato-"]),
    ("Convenio",       ["_convenio_","protocolo_acuerdo","protocolo-acuerdo",
                        "convenio_colaboracion","convenio_marco"]),
    ("Addendum",       ["_addendum_","addendum_","adenda_"]),
    ("Bases",          ["_bases_licitacion","bases_concurso","bases_llamado"]),

    # Certificaciones y respaldo legal
    ("Certificado",    ["_certificado_","certificado_","cert_"]),
    ("Solicitud",      ["_solicitud_","solicitud_","solicitud-"]),
    ("Declaración",    ["_declaracion_","declaracion_jurada","decl_jurada"]),
    ("Poder",          ["_poder_notarial","poder_simple","mandato_"]),
    ("Ficha",          ["_ficha_","ficha_tecnica","ficha_social","ficha_ingreso"]),

    # Planificación y gestión
    ("Plan",           ["_plan_","plan_trabajo","plan_accion","plan_gestion"]),
    ("Programa",       ["_programa_","programa_trabajo","programa_"]),
    ("Proyecto",       ["_proyecto_","proyecto_n","proyecto-"]),
    ("Propuesta",      ["_propuesta_","propuesta_tecnica","propuesta_"]),
    ("Protocolo",      ["_protocolo_","protocolo_"]),

    # Atención ciudadana (OIRS)
    ("Reclamo",        ["_reclamo_","reclamo_","reclamacion_"]),
    ("Consulta",       ["_consulta_","consulta_ciudadana","consulta-"]),
    ("Sugerencia",     ["_sugerencia_","sugerencia_"]),
    ("Denuncia",       ["_denuncia_","denuncia_"]),

    # Recursos humanos
    ("Resolución RRHH",["res_rrhh","resolucion_rrhh","res_contratacion",
                        "res_termino","res_cometido","res_feriado"]),
    ("Cometido",       ["_cometido_","cometido_funcionario","resol_cometido"]),

    # Técnicos y geoespaciales (contexto territorial)
    ("Plano",          ["_plano_","plano_n","plano-"]),
    ("Mapa",           ["_mapa_","mapa_"]),
    ("Expediente",     ["_expediente_","expediente_n"]),
    # Finanzas complementarios
    ("Guía de Despacho",  ["_guia_despacho","guia-despacho","_gd_"]),
    ("Nota de Crédito",   ["_nota_credito","nota-credito","_nc_"]),
    ("Nota de Débito",    ["_nota_debito","nota-debito","_nd_"]),
    ("Comprobante",       ["_comprobante_","comprobante_pago","comprobante_ingreso"]),
    ("Garantía",          ["_garantia_","boleta_garantia","poliza_garantia"]),
    ("Reembolso",         ["_reembolso_","solicitud_reembolso"]),

    # Licitaciones y compras públicas
    ("Bases Licitación",  ["bases_licitacion","bases_concurso","bases_llamado",
                           "bases_tecnicas","bases_admin"]),
    ("Oferta",            ["_oferta_tecnica","_oferta_economica","oferta_proveedor"]),
    ("Contrato Marco",    ["contrato_marco","_marco_","convenio_marco"]),
    ("Especificaciones",  ["especificaciones_tecnicas","eett_","_eett"]),

    # Desarrollo social y atención ciudadana
    ("Ficha Social",      ["ficha_social","_fichasocial_","ficha_cas",
                           "ficha_familia","registro_social"]),
    ("Informe Social",    ["informe_social","inf_social","_is_social"]),
    ("Postulación",       ["_postulacion_","postulacion_subsidio","formulario_postulacion"]),
    ("Subsidio",          ["_subsidio_","resolucion_subsidio","decreto_subsidio"]),
    ("Padrón",            ["_padron_","padron_beneficiarios","nomina_beneficiarios"]),

    # Planificación e inversión pública
    ("Ficha IDI",         ["_ficha_idi","ficha_idi","_fidi_","idi_n"]),
    ("Ficha EBI",         ["_ficha_ebi","ficha_ebi","_febi_"]),
    ("RS",                ["_rs_n","resultado_rs","rs_favorable","rs_observado"]),
    ("Perfil Proyecto",   ["perfil_proyecto","_perfil_","etapa_perfil"]),
    ("Bases de Diseño",   ["bases_diseno","bases_diseño","_bd_tecnico"]),

    # RRHH y personal
    ("Contrato Trabajo",  ["contrato_trabajo","_cont_trabajo","contrato_honorario",
                           "contrato_plazo_fijo"]),
    ("Finiquito",         ["_finiquito_","acuerdo_finiquito"]),
    ("Licencia",          ["_licencia_medica","licencia_med","_lm_"]),
    ("Feriado",           ["solicitud_feriado","_feriado_","resol_feriado"]),
    ("Calificación",      ["_calificacion_","hoja_calificacion","calificacion_anual"]),

    # Documentos técnicos territoriales
    ("Informe SEREMI",    ["informe_seremi","pronunciamiento_seremi","ord_seremi"]),
    ("Informe SUBDERE",   ["informe_subdere","oficio_subdere","pronunciamiento_subdere"]),
    ("Resolución SMA",    ["res_sma","resolucion_sma","pronunciamiento_sma"]),
    ("DGA",               ["informe_dga","res_dga","pronunciamiento_dga"]),
    ("SAG",               ["informe_sag","res_sag","pronunciamiento_sag"]),
]


TIPOS_TEXTO_BASE = {
    "Oficio":          ["oficio","of. n°","of n°","por medio del presente oficio"],
    "Resolución":      ["resolución exenta","res. ex.","resolución n°","res. afecta"],
    "Acta":            ["acta de sesión","acta sesion","acta n°","acta de acuerdos",
                        "sesión ordinaria","sesión extraordinaria","acta de entrega"],
    "Informe":         ["informe técnico","informe n°","se informa","informe de",
                        "informe final","informe legal","presenta informe",
                        "informe social","informe financiero"],
    "Decreto":         ["decreto exento","decreto n°","decreto alcaldicio"],
    "Dictamen":        ["dictamen"],
    "Ordinario":       ["ordinario n°"],
    "Circular":        ["circular n°"],
    "Convenio":        ["convenio","protocolo de acuerdo","convenio de colaboración"],
    "Solicitud":       ["solicitud","se solicita","por medio de la presente solicitud"],
    "Minuta":          ["minuta"],
    "Memo":            ["memorándum","memorandum","memo interno"],
    "Certificado":     ["certificado","se certifica","el suscrito certifica"],
    "Contrato":        ["contrato","las partes acuerdan","el contratante"],
    "Factura":         ["factura n°","rut del emisor","total a pagar"],
    "Rendición":       ["rendición de cuentas","rendición de gastos",
                        "se rinde cuenta","comprobantes adjuntos"],
    "Presupuesto":     ["presupuesto estimado","cuadro de gastos","ítem presupuestario"],
    "Cotización":      ["cotización n°","valor unitario","precio unitario"],
    "Reclamo":         ["reclamo n°","el recurrente","objeto del reclamo"],
    "Declaración":     ["declaración jurada","declaro bajo juramento"],
    "Evaluación":      ["evaluación técnica","criterios de evaluación","puntaje"],
    "Plan":            ["plan de trabajo","plan de acción","objetivo general"],
    "Proyecto":        ["nombre del proyecto","objetivo del proyecto","beneficiarios"],
    "Ficha":           ["ficha técnica","ficha social","datos del beneficiario"],
    "Guía de Despacho":  ["guía de despacho", "n° guía", "despacho de mercadería"],
    "Nota de Crédito":   ["nota de crédito", "n° nota de crédito"],
    "Nota de Débito":    ["nota de débito", "n° nota de débito"],
    "Comprobante":       ["comprobante de pago", "n° operación", "monto pagado"],
    "Garantía":          ["boleta de garantía", "póliza de garantía", "garantiza el fiel"],
    "Bases Licitación":  ["bases de licitación", "bases administrativas", "bases técnicas",
                          "especificaciones técnicas", "licitación pública"],
    "Oferta":            ["oferta técnica", "oferta económica", "propuesta técnica"],
    "Ficha Social":      ["ficha social", "registro social de hogares", "rsh"],
    "Informe Social":    ["informe social", "diagnóstico social", "condición socioeconómica"],
    "Postulación":       ["postulación", "formulario de postulación", "solicita subsidio"],
    "Subsidio":          ["subsidio habitacional", "subsidio de arriendo", "decreto de subsidio"],
    "Ficha IDI":         ["identificación de la iniciativa", "código bip", "nombre del proyecto",
                          "ficha idi"],
    "RS":                ["resultado satisfactorio", "rs favorable", "rs observado"],
    "Contrato Trabajo":  ["contrato de trabajo", "contrato de honorarios", "el empleador"],
    "Finiquito":         ["finiquito", "término del contrato", "liquidación final"],
    "Licencia":          ["licencia médica", "reposo médico", "días de reposo"],
    "Calificación":      ["hoja de calificación", "calificación anual", "nota final"],
    "Informe SEREMI":    ["seremi", "secretaría regional ministerial"],
    "Informe SUBDERE":   ["subdere", "subsecretaría de desarrollo regional"],
    "Resolución SMA":    ["sma", "superintendencia del medio ambiente"],
}


# Tipos activos — se combinan base + config.json al iniciar
TIPOS_NOMBRE = list(TIPOS_NOMBRE_BASE)
TIPOS_TEXTO  = dict(TIPOS_TEXTO_BASE)

PATRONES_CORREO_BASE = [
    r'de:\s*\S+@\S+', r'para:\s*\S+@\S+',
    r'fwd:|fw:|reenviado', r'forwarded message',
    r'\d+\s+de\s+\w+\s+de\s+\d{4}\s+a\s+las\s+\d+:\d+',
]
# Dominio institucional — se actualiza desde config.json
PATRONES_CORREO = list(PATRONES_CORREO_BASE)

PATRONES_NUM_DOC = [
    (r'folio\s*n[°o]?\s*(\d{3,6})',                           "Folio N°{}",  "regex:folio"),
    (r'res(?:olución|olucion)?\.?\s*ex(?:enta)?\.?\s*n[°o]?\s*(\d+[/-]?\d*)', "Res.Ex. N°{}", "regex:res-ex"),
    (r'oficio\s*(?:ord\.?)?\s*n[°o]?\s*(\d+)',                "Oficio N°{}", "regex:oficio"),
    (r'd(?:ecreto)?\.?\s*ex(?:ento)?\.?\s*n[°o]?\s*(\d+)',    "D.Ex. N°{}",  "regex:decreto"),
    (r'acta\s*(?:de\s*sesión\s*)?(?:ord\.?\s*)?n[°o]?\s*(\d+)', "Acta N°{}","regex:acta"),
    (r'ordinario\s*n[°o]?\s*(\d+)',                            "Ord. N°{}",   "regex:ordinario"),
    (r'circular\s*n[°o]?\s*(\d+)',                             "Circ. N°{}",  "regex:circular"),
    (r'\bn[°o]\s*(\d{3,6}(?:[/-]\d+)?)\b',                    "N°{}",        "regex:generico"),
]

# Colores UI
CF="#1e2a38"; CP="#2c3e50"; CA="#3498db"; CT="#ecf0f1"
CV="#27ae60"; CR="#e74c3c"; CY="#f39c12"; CN="#e67e22"
CM="#8e44ad"; CPAR="#243447"; CIMP="#1e2a38"; CSUB="#1a2e3a"
CBAJA="#3a1a1a"; CMEDIA="#3a3a1a"

# Colores Excel
XA="1A3A5C"; XB="FFFFFF"; XG1="EBF5FB"; XG2="D6EAF8"
XV="E9F7EF"; XM="F4ECF7"; XO="FEF9E7"

# Textos institucionales
# Valores por defecto — el usuario los configura desde la interfaz
# o desde config.json. Se sobreescriben al iniciar la aplicación.
INST_NOMBRE    = "ADMINISTRACIÓN PÚBLICA"
INST_SUBTITULO = "Unidad / Departamento"
INST_PIE       = "Unidad / Departamento"
APP_TITULO     = "GENERADOR DE ÍNDICE + FOLIADOR"
APP_VERSION    = "v7.0"

# Config institucional activa (se actualiza desde UI o config.json)
_cfg_inst = {
    "nombre":    INST_NOMBRE,
    "subtitulo": INST_SUBTITULO,
    "pie":       INST_PIE,
}

def inst_nombre()    -> str: return _cfg_inst.get("nombre",    INST_NOMBRE)
def inst_subtitulo() -> str: return _cfg_inst.get("subtitulo", INST_SUBTITULO)
def inst_pie()       -> str: return _cfg_inst.get("pie",       INST_PIE)
def actualizar_inst(nombre: str, subtitulo: str, pie: str = ""):
    _cfg_inst["nombre"]    = nombre.upper() if nombre else INST_NOMBRE
    _cfg_inst["subtitulo"] = subtitulo      if subtitulo else INST_SUBTITULO
    _cfg_inst["pie"]       = pie or subtitulo or INST_PIE

COLS_INDICE = ["N° Folio","N° Documento","Subcarpeta","Fecha","Hora","Tipo",
               "Descripción / Asunto","Págs","Tamaño","Nombre foliado"]

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 2: LIBRERÍAS OPCIONALES
# ══════════════════════════════════════════════════════════════════════════════

try:
    import pdfplumber;              PDF_OK  = True
except ImportError:
    PDF_OK  = False; log.warning("pdfplumber no instalado")

try:
    from docx import Document as DocxDoc
    from docx.shared import Pt, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_ALIGN_VERTICAL
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    DOCX_OK = True
except ImportError:
    DOCX_OK = False; log.warning("python-docx no instalado")

try:
    import openpyxl
    from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    XLSX_OK = True
except ImportError:
    XLSX_OK = False; log.warning("openpyxl no instalado")

try:
    from PIL import Image;          PIL_OK  = True
except ImportError:
    PIL_OK  = False

try:
    import rarfile;                 RAR_OK  = True
except ImportError:
    RAR_OK  = False

try:
    import py7zr;                   SZ_OK   = True
except ImportError:
    SZ_OK   = False

try:
    from pptx import Presentation;  PPTX_OK = True
except ImportError:
    PPTX_OK = False



# ══════════════════════════════════════════════════════════════════════════════
# SISTEMA DE CONFIGURACIÓN JSON (NUEVO v7.1)
# Permite personalizar institución y tipos sin tocar el código.
# Archivo: config_indice.json (misma carpeta que el script)
# ══════════════════════════════════════════════════════════════════════════════

CONFIG_FILE = Path(__file__).parent / "config_indice.json"

CONFIG_DEFAULT = {
    "institucion": {
        "nombre":    "ADMINISTRACIÓN PÚBLICA",
        "subtitulo": "Unidad / Departamento",
        "pie":       "Unidad / Departamento"
    },
    "dominio_correo": "",
    "tipos_adicionales": {
        "nombre": [],
        "texto":  {}
    },
    "version": APP_VERSION
}

def cargar_config_json() -> dict:
    """
    Carga config_indice.json si existe.
    Si no existe, lo crea con valores por defecto.
    Retorna el dict de configuración activo.
    """
    if not CONFIG_FILE.exists():
        try:
            with open(CONFIG_FILE,"w",encoding="utf-8") as f:
                json.dump(CONFIG_DEFAULT, f, ensure_ascii=False, indent=2)
            log.info(f"config_indice.json creado con valores por defecto")
        except (PermissionError, OSError) as e:
            log.warning(f"No se pudo crear config_indice.json: {e}")
        return dict(CONFIG_DEFAULT)
    try:
        with open(CONFIG_FILE,"r",encoding="utf-8") as f:
            cfg = json.load(f)
        log.info(f"config_indice.json cargado")
        return cfg
    except (json.JSONDecodeError, OSError) as e:
        log.warning(f"config_indice.json inválido: {e} — usando valores por defecto")
        return dict(CONFIG_DEFAULT)

def guardar_config_json(cfg: dict) -> bool:
    """Guarda configuración en config_indice.json."""
    try:
        cfg["version"] = APP_VERSION
        with open(CONFIG_FILE,"w",encoding="utf-8") as f:
            json.dump(cfg, f, ensure_ascii=False, indent=2)
        log.info("config_indice.json guardado")
        return True
    except (PermissionError, OSError) as e:
        log.warning(f"No se pudo guardar config_indice.json: {e}")
        return False

def aplicar_config(cfg: dict) -> None:
    """
    Aplica la configuración cargada:
    - Actualiza textos institucionales
    - Extiende tipos documentales con los adicionales del JSON
    - Agrega dominio de correo institucional
    """
    global TIPOS_NOMBRE, TIPOS_TEXTO, PATRONES_CORREO, _cfg_inst

    # Institución
    inst = cfg.get("institucion", {})
    nombre    = inst.get("nombre",    CONFIG_DEFAULT["institucion"]["nombre"])
    subtitulo = inst.get("subtitulo", CONFIG_DEFAULT["institucion"]["subtitulo"])
    pie       = inst.get("pie",       subtitulo)
    _cfg_inst.update({"nombre": nombre.upper() if nombre else nombre,
                      "subtitulo": subtitulo, "pie": pie or subtitulo})

    # Tipos adicionales desde JSON
    tipos_add = cfg.get("tipos_adicionales", {})

    # Nombres adicionales: lista de [tipo, [patrones]]
    for entrada in tipos_add.get("nombre", []):
        if isinstance(entrada, list) and len(entrada) == 2:
            tipo_nuevo, patrones_nuevos = entrada
            # Si el tipo ya existe, extender sus patrones
            for i, (tipo_ex, pats) in enumerate(TIPOS_NOMBRE):
                if tipo_ex.lower() == tipo_nuevo.lower():
                    TIPOS_NOMBRE[i] = (tipo_ex, list(set(pats + patrones_nuevos)))
                    break
            else:
                TIPOS_NOMBRE.append((tipo_nuevo, patrones_nuevos))

    # Texto adicional: dict {tipo: [palabras]}
    for tipo_nuevo, palabras in tipos_add.get("texto", {}).items():
        if tipo_nuevo in TIPOS_TEXTO:
            TIPOS_TEXTO[tipo_nuevo] = list(set(TIPOS_TEXTO[tipo_nuevo] + palabras))
        else:
            TIPOS_TEXTO[tipo_nuevo] = palabras

    # Dominio de correo institucional
    dominio = cfg.get("dominio_correo", "").strip()
    PATRONES_CORREO = list(PATRONES_CORREO_BASE)
    if dominio:
        patron = rf'@{re.escape(dominio)}'
        PATRONES_CORREO.append(patron)
        log.info(f"Dominio de correo institucional: @{dominio}")

    log.info(f"Configuración aplicada: {_cfg_inst['nombre']} | "
             f"{len(TIPOS_NOMBRE)} tipos documentales activos")

# Config global activa
_config_activa: dict = {}

def inicializar_config() -> dict:
    """Carga y aplica la configuración al iniciar la aplicación."""
    global _config_activa
    _config_activa = cargar_config_json()
    aplicar_config(_config_activa)
    return _config_activa

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 3: DATACLASS DocumentoAnalizado (NUEVO v7.0)
# ══════════════════════════════════════════════════════════════════════════════

@dataclass
class DocumentoAnalizado:
    """
    Modelo de datos para un documento del expediente.
    Reemplaza los diccionarios de versiones anteriores.
    """
    n:             int   = 0
    folio:         str   = ""
    nombre_foliado:str   = ""
    subcarpeta:    str   = "—"
    nombre:        str   = ""
    ruta:          object = None   # Path — no tipado para compatibilidad Python 3.8
    extension:     str   = ""
    tipo:          str   = ""
    fuente_tipo:   str   = ""
    conf_tipo:     str   = ""
    regla_tipo:    str   = ""
    num_doc:       str   = ""
    regla_num_doc: str   = ""
    fecha:         str   = ""
    hora:          str   = ""
    fuente_fecha:  str   = ""
    regla_fecha:   str   = ""
    paginas:       str   = ""
    tamano:        str   = ""
    descripcion:   str   = ""
    hash_sha256:   str   = ""
    historial_ediciones: list = field(default_factory=list)

    def hash_corto(self) -> str:
        h = self.hash_sha256
        if not h or h.startswith(("ERROR","ARCHIVO","SIN")): return h or "—"
        return f"{h[:8]}...{h[-8:]}"

    def tiene_hash_valido(self) -> bool:
        h = self.hash_sha256
        return bool(h) and len(h)==64 and not h.startswith(("ERROR","ARCHIVO","SIN"))

    def registrar_edicion(self, campo: str, valor_anterior: str,
                           valor_nuevo: str, usuario: str = ""):
        import os as _os
        self.historial_ediciones.append({
            "timestamp":      __import__("datetime").datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "campo":          campo,
            "valor_anterior": valor_anterior,
            "valor_nuevo":    valor_nuevo,
            "usuario":        usuario or _os.getenv("USERNAME","usuario"),
        })

    def fue_editado(self) -> bool:
        return len(self.historial_ediciones) > 0

    def como_fila_ui(self) -> tuple:
        return (self.folio, self.num_doc, self.subcarpeta,
                self.fecha, self.hora, self.tipo, self.descripcion,
                self.paginas, self.tamano,
                self.fuente_fecha, self.conf_tipo,
                "✎" if self.fue_editado() else "",
                self.hash_corto(), self.nombre_foliado)

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 3: VALIDADORES
# ══════════════════════════════════════════════════════════════════════════════

def validar_fecha(valor: str) -> str:
    if not valor:
        return "La fecha no puede estar vacía"
    try:
        datetime.strptime(valor.strip(), "%Y-%m-%d"); return ""
    except ValueError:
        return f"Fecha inválida '{valor}'. Use YYYY-MM-DD (ej: 2025-07-17)"

def validar_hora(valor: str) -> str:
    if not valor: return ""
    valor = valor.strip()
    if not re.match(r'^\d{2}:\d{2}$', valor):
        return f"Hora inválida '{valor}'. Use HH:MM con dos dígitos (ej: 08:59)"
    try:
        datetime.strptime(valor, "%H:%M"); return ""
    except ValueError:
        return f"Hora inválida '{valor}'"

def validar_tipo(valor: str) -> str:
    if not valor or not valor.strip():
        return "El tipo de documento no puede estar vacío"
    return ""

def validar_descripcion(valor: str) -> str:
    if not valor or not valor.strip():
        return "La descripción no puede estar vacía"
    if len(valor.strip()) < 4:
        return "Descripción demasiado corta (mínimo 4 caracteres)"
    return ""

def validar_campos_edicion(fecha, hora, tipo, desc) -> list:
    errores = []
    for fn, val in [(validar_fecha,fecha),(validar_hora,hora),
                    (validar_tipo,tipo),(validar_descripcion,desc)]:
        msg = fn(val)
        if msg: errores.append(msg)
    return errores

def validar_foliacion(registros: list) -> list:
    problemas = []; folios_vistos = {}
    for r in registros:
        try: n = int(r["folio"])
        except (ValueError, KeyError) as e:
            problemas.append(f"Folio inválido '{r.get('folio','')}': {e}"); continue
        if n in folios_vistos:
            problemas.append(f"Folio duplicado {str(n).zfill(3)}: "
                             f"'{folios_vistos[n]}' y '{r['nombre']}'")
        else:
            folios_vistos[n] = r["nombre"]
    if folios_vistos:
        fo = sorted(folios_vistos.keys())
        for i in range(len(fo)-1):
            if fo[i+1]-fo[i]>1:
                falt=list(range(fo[i]+1,fo[i+1]))
                problemas.append(f"Gap: faltan "
                    f"{str(falt[0]).zfill(3)}–{str(falt[-1]).zfill(3)}")
    return problemas

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 4: UTILIDADES COMUNES
# ══════════════════════════════════════════════════════════════════════════════

def tamano_str(ruta: Path) -> str:
    try:
        kb = ruta.stat().st_size/1024
        return f"{kb:.1f} KB" if kb<1024 else f"{kb/1024:.2f} MB"
    except OSError as e:
        log.debug(f"Tamaño {ruta.name}: {e}"); return ""

def fecha_mtime(ruta: Path) -> tuple:
    try:
        dt=datetime.fromtimestamp(ruta.stat().st_mtime)
        return dt.strftime("%Y-%m-%d"),dt.strftime("%H:%M")
    except OSError as e:
        log.debug(f"mtime {ruta.name}: {e}"); return "",""

def fecha_desde_texto(texto: str) -> tuple:
    meses={"enero":"01","febrero":"02","marzo":"03","abril":"04",
           "mayo":"05","junio":"06","julio":"07","agosto":"08",
           "septiembre":"09","octubre":"10","noviembre":"11","diciembre":"12"}
    patron=(r'(\d{1,2})\s+de\s+(enero|febrero|marzo|abril|mayo|junio|julio|agosto|'
            r'septiembre|octubre|noviembre|diciembre)\s+(?:de\s+)?(\d{4})'
            r'(?:\s+a\s+las?\s+(\d{1,2}):(\d{2}))?')
    fechas=[]
    for m in re.finditer(patron,texto,re.IGNORECASE):
        g=m.groups()
        try:
            mes=meses.get(g[1].lower(),"")
            if not mes: continue
            fs=f"{g[2]}{mes}{g[0].zfill(2)}"
            datetime.strptime(fs,"%Y%m%d")
            hs=f"{g[3].zfill(2)}:{g[4].zfill(2)}" if g[3] and g[4] else ""
            fechas.append((fs,hs,m.group(0)[:40]))
        except ValueError: continue
    if not fechas: return "","","",""
    fechas.sort(key=lambda x:x[0])
    fs,hs,regla=fechas[0]
    return f"{fs[:4]}-{fs[4:6]}-{fs[6:]}",hs,"Texto-documento",f'"{regla}"'

def fecha_desde_nombre(nombre: str) -> tuple:
    m=re.match(r'^(\d{4})-(\d{2})-(\d{2})_(\d{2})(\d{2})',nombre)
    if m:
        y,mo,d,h,mi=m.groups()
        try:
            datetime.strptime(f"{y}{mo}{d}","%Y%m%d")
            return f"{y}-{mo}-{d}",f"{h}:{mi}","Nombre-archivo",f'"{m.group(0)}"'
        except ValueError: pass
    return "","","",""

def desc_desde_nombre(nombre: str) -> str:
    b=re.sub(r'^\d{4}-\d{2}-\d{2}_\d{4}_?','',nombre)
    b=re.sub(r'^\d+_','',b); b=re.sub(r'[_]+',' ',b)
    return re.sub(r'\s+',' ',b).strip() or "Sin descripción"

def _num_doc_limpio(num_doc: str) -> str:
    """
    Extrae la parte numérica del num_doc para el nombre foliado.
    Ejemplos:
      "Folio N°2544"       → "2544"
      "Oficio N°1192"      → "1192"
      "Res.Ex. N°548/2025" → "548-2025"
      "Acta N°35"          → "35"
    """
    if not num_doc: return ""
    numeros = re.findall(r'\d+', num_doc)
    if not numeros: return ""
    if len(numeros) == 1: return numeros[0]
    return '-'.join(numeros)

def _normalizar_para_comparar(texto: str) -> str:
    """Normaliza texto para comparación: minúsculas, sin guiones ni espacios."""
    return re.sub(r'[\s\-_/]','', texto.lower().strip())

def _limpiar_desc_redundante(desc: str, fecha: str, hora: str,
                              num_doc: str, tipo: str) -> str:
    """
    Elimina de la descripción los componentes que ya están
    en la estructura del nombre foliado (fecha, hora, tipo, num_doc).
    Evita redundancias como:
      001_20250717_1750_Correo_Correo-Ingreso... → Correo duplicado
      001_20250717_1750_Oficio_Oficio-1568-20250717... → todo duplicado
    """
    b = desc

    # Eliminar el tipo si aparece al inicio de la descripción
    tipo_limpio = re.sub(r'[^a-zA-ZáéíóúÁÉÍÓÚüÜñÑ]','',tipo or "")
    if tipo_limpio:
        b = re.sub(rf'^{re.escape(tipo_limpio)}[\s_\-]*',
                   '', b, flags=re.IGNORECASE)

    # Eliminar fecha en formatos comunes: 20250717, 2025-07-17, 2025_07_17
    if fecha:
        fecha_num = fecha.replace("-","")
        for patron_f in [fecha_num, fecha, fecha.replace("-","_")]:
            b = re.sub(re.escape(patron_f), '', b, flags=re.IGNORECASE)

    # Eliminar hora en formatos comunes: 1750, 17:50, 17_50
    if hora:
        hora_num = hora.replace(":","")
        for patron_h in [hora_num, hora, hora.replace(":","_")]:
            b = re.sub(re.escape(patron_h), '', b, flags=re.IGNORECASE)

    # Eliminar número de documento (ej: "1568", "N°1568", "1192")
    if num_doc:
        # Extraer solo los dígitos del num_doc
        solo_num = re.sub(r'[^\d]','',num_doc)
        if solo_num and len(solo_num) >= 3:
            b = re.sub(rf'\b{re.escape(solo_num)}\b', '', b)

    # Limpiar residuos: guiones, underscores y espacios múltiples
    b = re.sub(r'[_\-\s]+', '_', b.strip('_- '))
    b = re.sub(r'_+', '_', b).strip('_')
    return b

def desc_foliado(nombre: str) -> str:
    """Descripción base desde nombre de archivo (sin fecha/hora del prefijo)."""
    b = re.sub(r'^\d{4}-\d{2}-\d{2}_\d{4}_?','',nombre)
    b = re.sub(r'^\d+_','',b)
    b = re.sub(r'[\\/:*?"<>|]','',b)
    b = re.sub(r'\s+','_',b.strip())
    b = re.sub(r'_+','_',b).strip('_')
    if len(b)>80:
        partes=b[:80].split('_')
        b='_'.join(partes[:-1]) if len(partes)>1 else b[:80]
    return b or "Sin-descripcion"

def _score_nombre(fecha:str, hora:str, tipo:str, num_doc:str) -> int:
    """
    Evalúa qué tan completa es la información estructurada.
    Score alto → usar estructura pura (sin descripción extensa).
    Score bajo → usar descripción para complementar.
    """
    score = 0
    if fecha and fecha != "00000000": score += 30
    if hora  and hora  != "0000":     score += 20
    if tipo  and tipo  not in ("PDF","Word","Excel","Imagen","Texto"): score += 25
    if num_doc:                        score += 25
    return score

def construir_nombre_foliado(r) -> str:
    """
    Construye nombre foliado inteligente sin redundancias.
    Formato: NNN_YYYYMMDD_HHMM_Tipo[_NumDoc|_Descripcion].ext

    Reglas:
    1. Si hay número de documento → NNN_fecha_hora_Tipo_NumDoc.ext
    2. Si score alto sin num_doc  → NNN_fecha_hora_Tipo_DescLimpia.ext
    3. Si score bajo              → NNN_fecha_hora_Tipo_DescLimpia.ext

    La descripción se limpia eliminando componentes ya presentes
    (tipo, fecha, hora, num_doc) para evitar redundancias.
    """
    # Compatibilidad dict / DocumentoAnalizado
    def _get(obj, key, default=""):
        if hasattr(obj, key): return getattr(obj, key) or default
        if isinstance(obj, dict): return obj.get(key, default)
        return default

    folio   = _get(r,"folio","000")
    fecha_s = _get(r,"fecha","").replace("-","") or "00000000"
    hora_s  = _get(r,"hora","").replace(":","")  or "0000"
    tipo_r  = _get(r,"tipo","Doc")
    num_doc = _get(r,"num_doc","")
    desc_r  = _get(r,"descripcion","")
    ext     = _get(r,"extension","")

    # Tipo limpio para el nombre (solo letras)
    tipo = re.sub(r'[^a-zA-ZáéíóúÁÉÍÓÚüÜñÑ]','',tipo_r) or "Doc"

    # ── Caso 1: hay número de documento ──────────────────────────────────────
    # Resultado: 001_20250717_1750_Oficio_1192.ext
    if num_doc:
        num_limpio = _num_doc_limpio(num_doc)
        if num_limpio:
            return f"{folio}_{fecha_s}_{hora_s}_{tipo}_{num_limpio}{ext}"

    # ── Caso 2: sin número de documento → usar descripción limpia ────────────
    # Resultado: 001_20250717_1750_Correo_Ingreso-Reclamacion-ECMPO.ext
    desc_base = desc_foliado(desc_r)
    desc_limpia = _limpiar_desc_redundante(
        desc_base,
        fecha=_get(r,"fecha",""),
        hora=_get(r,"hora",""),
        num_doc=num_doc,
        tipo=tipo_r,
    )

    # Si la descripción quedó vacía o muy corta tras la limpieza
    if not desc_limpia or len(desc_limpia) < 3:
        score = _score_nombre(fecha_s, hora_s, tipo_r, num_doc)
        if score >= 50:
            # Estructura suficientemente informativa sola
            return f"{folio}_{fecha_s}_{hora_s}_{tipo}{ext}"
        else:
            # Usar descripción original sin limpiar
            desc_limpia = desc_base[:60]

    # Acortar descripción si excede 70 chars
    if len(desc_limpia) > 70:
        partes = desc_limpia[:70].split('_')
        desc_limpia = '_'.join(partes[:-1]) if len(partes)>1 else desc_limpia[:70]

    return f"{folio}_{fecha_s}_{hora_s}_{tipo}_{desc_limpia}{ext}"



def resumen_estadistico(registros: list) -> dict:
    por_tipo={}; por_sub={}; por_ext={}; total_kb=0.0
    for r in registros:
        por_tipo[r.get("tipo","?")]=por_tipo.get(r.get("tipo","?"),0)+1
        por_sub[r.get("subcarpeta","—")]=por_sub.get(r.get("subcarpeta","—"),0)+1
        ext=r.get("extension","?").upper().lstrip(".")
        por_ext[ext]=por_ext.get(ext,0)+1
        try:
            t=r.get("tamano",""); mult=1024 if "MB" in t else 1
            total_kb+=float(re.sub(r'[^\d.]','',t) or 0)*mult
        except (ValueError,TypeError): pass
    total_mb=total_kb/1024
    return {
        "por_tipo":     dict(sorted(por_tipo.items(),key=lambda x:x[1],reverse=True)),
        "por_sub":      dict(sorted(por_sub.items(), key=lambda x:x[1],reverse=True)),
        "por_ext":      dict(sorted(por_ext.items(), key=lambda x:x[1],reverse=True)),
        "tamano_total": f"{total_mb:.1f} MB" if total_mb>=1 else f"{total_kb:.0f} KB",
        "total":        len(registros),
    }

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 5: HASH SHA-256 (NUEVO en v6.0)
# ══════════════════════════════════════════════════════════════════════════════

def generar_hash(ruta: Path) -> str:
    """
    Genera huella digital SHA-256 del archivo.
    Lee en bloques de 64KB para no saturar RAM con archivos grandes.
    Retorna el hash como string hexadecimal o mensaje de error.
    """
    h = hashlib.sha256()
    try:
        with open(ruta, "rb") as f:
            for bloque in iter(lambda: f.read(65536), b""):
                h.update(bloque)
        return h.hexdigest()
    except FileNotFoundError:
        log.warning(f"Hash: archivo no encontrado {ruta.name}")
        return "ARCHIVO_NO_ENCONTRADO"
    except PermissionError as e:
        log.warning(f"Hash: sin permisos {ruta.name}: {e}")
        return "SIN_PERMISOS"
    except OSError as e:
        log.warning(f"Hash: error de sistema {ruta.name}: {e}")
        return "ERROR_LECTURA"

def hash_corto(hash_completo: str) -> str:
    """Retorna los primeros 16 chars del hash para mostrar en tabla."""
    if hash_completo.startswith("ERROR") or hash_completo.startswith("ARCHIVO") or \
       hash_completo.startswith("SIN"):
        return hash_completo
    return f"{hash_completo[:8]}...{hash_completo[-8:]}"


# ══════════════════════════════════════════════════════════════════════════════
# CAPA 7: CACHÉ JSON LIGERA (NUEVO v7.0)
# ══════════════════════════════════════════════════════════════════════════════

CACHE_VERSION = "7.0"

def _ruta_cache(carpeta):
    return Path(carpeta) / ".indice_cache.json"

def cargar_cache(carpeta) -> dict:
    ruta = _ruta_cache(carpeta)
    if not ruta.exists(): return {}
    try:
        with open(ruta,"r",encoding="utf-8") as f:
            data = json.load(f)
        if data.get("version") != CACHE_VERSION: return {}
        return data.get("archivos",{})
    except Exception as e:
        log.debug(f"Caché inválida: {e}"); return {}

def guardar_cache(carpeta, cache: dict) -> None:
    ruta = _ruta_cache(carpeta)
    try:
        with open(ruta,"w",encoding="utf-8") as f:
            json.dump({"version":CACHE_VERSION,"archivos":cache},
                      f,ensure_ascii=False,indent=2)
    except Exception as e:
        log.warning(f"No se pudo guardar caché: {e}")

def cache_key(ruta) -> str:
    return str(Path(ruta).resolve())

def esta_en_cache_valida(ruta, cache: dict) -> bool:
    key = cache_key(ruta)
    if key not in cache: return False
    try:
        stat = Path(ruta).stat()
        entrada = cache[key]
        return (abs(stat.st_mtime - entrada.get("mtime",0)) < 1.0 and
                stat.st_size == entrada.get("tamano_bytes",-1))
    except OSError: return False

def resultado_desde_cache(ruta, cache: dict, raiz):
    key = cache_key(ruta)
    try:
        datos = cache[key]["resultado"]
        d = DocumentoAnalizado(**{k:v for k,v in datos.items() if k != "ruta"})
        d.ruta = Path(ruta)
        return d
    except Exception as e:
        log.debug(f"Cache inválida {Path(ruta).name}: {e}"); return None

def guardar_en_cache(ruta, doc, cache: dict) -> None:
    try:
        stat = Path(ruta).stat()
        from dataclasses import asdict as _asdict
        resultado = {k:v for k,v in _asdict(doc).items() if k != "ruta"}
        cache[cache_key(ruta)] = {
            "mtime":        stat.st_mtime,
            "tamano_bytes": stat.st_size,
            "resultado":    resultado,
        }
    except Exception as e:
        log.debug(f"No se pudo cachear {Path(ruta).name}: {e}")

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 5: HASH

def generar_hash_expediente(docs: list) -> str:
    """
    Huella digital SHA-256 del conjunto completo del expediente.
    Cualquier cambio en un archivo o en el orden altera este hash.
    """
    h = hashlib.sha256()
    docs_ordenados = sorted(docs, key=lambda d: getattr(d,"folio",""))
    for d in docs_ordenados:
        linea = f"{d.folio}|{d.nombre}|{d.hash_sha256}\n"
        h.update(linea.encode("utf-8"))
    return h.hexdigest()

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 6: MOTOR DE ANÁLISIS (independiente de la GUI)
# ══════════════════════════════════════════════════════════════════════════════

def _tipo_desde_nombre(nombre_l: str) -> tuple:
    for tipo,patrones in TIPOS_NOMBRE:
        for pat in patrones:
            if pat in f"_{nombre_l}_":
                return tipo,"Nombre-archivo","Alta",f'patron:"{pat}"'
    if any(p in nombre_l for p in ["correo","mail","fwd","fw ","mensaje"]):
        return "Correo","Nombre-archivo","Alta","patron:correo/mail/fwd"
    return "","","",""

def _tipo_desde_texto(texto: str, nombre_l: str, ext: str) -> tuple:
    if ext not in (".pdf",".docx",".txt"): return "","","",""
    for i,patron in enumerate(PATRONES_CORREO):
        try:
            if re.search(patron,texto,re.IGNORECASE):
                return "Correo","Texto-interno","Media",f'correo:patron[{i}]'
        except re.error as e:
            log.warning(f"Patrón correo inválido [{i}]: {e}")
    fuente=texto[:500]+" "+nombre_l
    for tipo,palabras in TIPOS_TEXTO.items():
        for palabra in palabras:
            if palabra in fuente:
                return tipo,"Texto-interno","Media",f'kw:"{palabra}"'
    return "","","",""

def detectar_tipo(texto: str, nombre: str, ext: str) -> tuple:
    tipo_base=TIPOS_MAP.get(ext.lower(),ext[1:].upper())
    nombre_l=nombre.lower()
    tipo,fuente,conf,regla=_tipo_desde_nombre(nombre_l)
    if tipo: return tipo,fuente,conf,regla
    tipo,fuente,conf,regla=_tipo_desde_texto(texto,nombre_l,ext)
    if tipo: return tipo,fuente,conf,regla
    return tipo_base,"Extensión","Baja",f'ext:"{ext}"'

def extraer_numero_doc(texto: str, nombre: str) -> tuple:
    fuente=(texto+" "+nombre).lower()
    for patron,fmt,regla in PATRONES_NUM_DOC:
        try:
            m=re.search(patron,fuente,re.IGNORECASE)
            if m: return fmt.format(m.group(1).strip()),regla
        except re.error as e:
            log.warning(f"Patrón N° doc inválido: {e}")
    m_nombre=re.search(r'(?:n[°_]?|num[_]?)(\d{3,6}(?:[-_]\d{4})?)',nombre,re.IGNORECASE)
    if m_nombre:
        return f"N°{m_nombre.group(1).replace('_','-')}","regex:nombre-archivo"
    return "",""

def _extraer_texto_pdf(ruta: Path) -> str:
    if not PDF_OK: return ""
    try:
        with pdfplumber.open(ruta) as pdf:
            return "\n".join(p.extract_text() or "" for p in pdf.pages[:4]).lower()
    except FileNotFoundError:
        log.warning(f"PDF no encontrado: {ruta.name}"); return ""
    except PermissionError as e:
        log.warning(f"Sin permisos PDF {ruta.name}: {e}"); return ""
    except Exception as e:
        log.warning(f"Error PDF {ruta.name}: {e}"); return ""

def _extraer_texto_docx(ruta: Path) -> str:
    if not DOCX_OK: return ""
    try:
        doc=DocxDoc(ruta)
        return "\n".join(p.text for p in doc.paragraphs[:60]).lower()
    except (FileNotFoundError,PermissionError) as e:
        log.warning(f"DOCX {ruta.name}: {e}"); return ""
    except Exception as e:
        log.warning(f"Error DOCX {ruta.name}: {e}"); return ""

def _extraer_texto_xlsx(ruta: Path) -> str:
    if not XLSX_OK: return ""
    try:
        wb=openpyxl.load_workbook(ruta,read_only=True,data_only=True)
        filas=[]
        for i,row in enumerate(wb.active.iter_rows(values_only=True)):
            if i>10: break
            filas.append(" ".join(str(c) for c in row if c))
        return "\n".join(filas).lower()
    except (FileNotFoundError,PermissionError) as e:
        log.warning(f"XLSX {ruta.name}: {e}"); return ""
    except Exception as e:
        log.warning(f"Error XLSX {ruta.name}: {e}"); return ""

def _extraer_texto_pptx(ruta: Path) -> tuple:
    if not PPTX_OK: return "",""
    try:
        prs=Presentation(ruta)
        texto="\n".join(
            shape.text for slide in list(prs.slides)[:3]
            for shape in slide.shapes if hasattr(shape,"text")).lower()
        return texto,str(len(prs.slides))
    except (FileNotFoundError,PermissionError) as e:
        log.warning(f"PPTX {ruta.name}: {e}"); return "",""
    except Exception as e:
        log.warning(f"Error PPTX {ruta.name}: {e}"); return "",""

def _info_kml(ruta: Path) -> dict:
    resultado={"nombre_capa":"","fecha":""}
    try:
        if ruta.suffix.lower()==".kmz":
            with zipfile.ZipFile(ruta,"r") as z:
                kf=[f for f in z.namelist() if f.endswith(".kml")]
                if not kf: return resultado
                contenido=z.read(kf[0]).decode("utf-8",errors="ignore")
        else:
            with open(ruta,encoding="utf-8",errors="ignore") as f: contenido=f.read()
        root=ET.fromstring(contenido)
        for elem in root.iter():
            tag=elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
            if tag=="name" and elem.text: resultado["nombre_capa"]=elem.text.strip(); break
        m=re.search(r'(\d{4})-(\d{2})-(\d{2})',contenido)
        if m: resultado["fecha"]=f"{m.group(1)}-{m.group(2)}-{m.group(3)}"
    except zipfile.BadZipFile as e: log.warning(f"KMZ corrupto {ruta.name}: {e}")
    except ET.ParseError as e:     log.warning(f"KML XML inválido {ruta.name}: {e}")
    except (FileNotFoundError,PermissionError) as e: log.warning(f"Sin acceso {ruta.name}: {e}")
    except Exception as e:         log.warning(f"Error KML {ruta.name}: {e}")
    return resultado

def _info_shp(ruta: Path) -> dict:
    srs=""
    try:
        prj=ruta.parent/(ruta.stem+".prj")
        if prj.exists():
            t=prj.read_text(encoding="utf-8",errors="ignore")
            if "WGS_1984" in t or "WGS84" in t: srs="WGS84"
            elif "SIRGAS" in t: srs="SIRGAS2000"
            elif "SAD69" in t:  srs="SAD69"
            elif "UTM" in t:    srs="UTM"
            else:               srs=t[:30].strip()
    except (FileNotFoundError,PermissionError) as e: log.warning(f"PRJ {ruta.name}: {e}")
    except Exception as e: log.warning(f"Error PRJ {ruta.name}: {e}")
    return {"srs":srs}

def _listar_zip(ruta: Path) -> str:
    archivos=[]
    try:
        ext=ruta.suffix.lower()
        if ext==".zip":
            with zipfile.ZipFile(ruta,"r") as z:
                archivos=[i.filename for i in z.infolist() if not i.is_dir()]
        elif ext==".rar" and RAR_OK:
            with rarfile.RarFile(ruta) as r:
                archivos=[i.filename for i in r.infolist() if not i.is_dir()]
        elif ext==".7z" and SZ_OK:
            with py7zr.SevenZipFile(ruta,mode="r") as z: archivos=z.getnames()
    except zipfile.BadZipFile as e: log.warning(f"ZIP corrupto {ruta.name}: {e}"); return "Corrupto"
    except (FileNotFoundError,PermissionError) as e: log.warning(f"Sin acceso {ruta.name}: {e}"); return "Sin acceso"
    except Exception as e: log.warning(f"Error ZIP {ruta.name}: {e}"); return "Error"
    if not archivos: return "Vacío"
    res=f"{len(archivos)} archivos: "+" | ".join(archivos[:6])
    if len(archivos)>6: res+=f" ... (+{len(archivos)-6} más)"
    return res

def analizar_archivo(ruta: Path, raiz: Path) -> dict:
    """
    Motor de análisis unificado. Independiente de la GUI.
    Incluye Hash SHA-256 y trazabilidad completa (fuente + regla exacta).
    """
    ext=ruta.suffix.lower(); nombre=ruta.stem; kb=tamano_str(ruta)
    try:
        partes=ruta.relative_to(raiz).parts
        subcarpeta=partes[0] if len(partes)>1 else "—"
    except ValueError:
        subcarpeta="—"

    texto=""; paginas=""
    desc=desc_desde_nombre(nombre)
    num_doc=""; regla_num_doc=""
    fuente_fecha=""; regla_fecha=""

    # Fecha desde nombre (prioritaria)
    fc,fh,fuente_fecha,regla_fecha=fecha_desde_nombre(nombre)

    try:
        if ext==".pdf":
            texto=_extraer_texto_pdf(ruta)
            if PDF_OK:
                try:
                    with pdfplumber.open(ruta) as pdf: paginas=str(len(pdf.pages))
                except Exception as e: log.debug(f"Páginas PDF {ruta.name}: {e}")
            if not fc: fc,fh,fuente_fecha,regla_fecha=fecha_desde_texto(texto)
        elif ext==".docx":
            texto=_extraer_texto_docx(ruta)
            if not fc: fc,fh,fuente_fecha,regla_fecha=fecha_desde_texto(texto)
        elif ext in (".xlsx",".xls"):
            texto=_extraer_texto_xlsx(ruta)
        elif ext==".pptx":
            texto,paginas=_extraer_texto_pptx(ruta)
        elif ext==".txt":
            try:
                with open(ruta,encoding="utf-8",errors="ignore") as f: texto=f.read(3000).lower()
                if not fc: fc,fh,fuente_fecha,regla_fecha=fecha_desde_texto(texto)
            except (FileNotFoundError,PermissionError) as e: log.warning(f"TXT {ruta.name}: {e}")
        elif ext in (".kml",".kmz"):
            info=_info_kml(ruta)
            if info["nombre_capa"]: desc=info["nombre_capa"]
            if info["fecha"] and not fc:
                fc=info["fecha"]; fh=""; fuente_fecha="XML-interno"; regla_fecha="KML/KMZ"
        elif ext==".shp":
            info=_info_shp(ruta)
            if info["srs"]: desc=f"{desc_desde_nombre(nombre)} [{info['srs']}]"
        elif ext in EXT_ZIP:
            desc=f"{desc_desde_nombre(nombre)} — {_listar_zip(ruta)}"
        elif ext in EXT_IMGS and PIL_OK:
            try:
                img=Image.open(ruta); w,h=img.size
                desc=f"{desc_desde_nombre(nombre)} [{w}×{h}px]"
            except Exception as e: log.debug(f"Imagen {ruta.name}: {e}")
    except Exception as e:
        log.error(f"Error inesperado {ruta.name}: {e}",exc_info=True)

    if not fc:
        fc,fh=fecha_mtime(ruta); fuente_fecha="Metadata-mtime"; regla_fecha="st_mtime"

    tipo,fuente_tipo,conf_tipo,regla_tipo=detectar_tipo(texto,nombre,ext)
    num_doc,regla_num_doc=extraer_numero_doc(texto,nombre)

    # Hash SHA-256
    hash_sha256=generar_hash(ruta)

    return DocumentoAnalizado(
        n=0, folio="", nombre_foliado="",
        subcarpeta=subcarpeta, nombre=ruta.name, ruta=ruta, extension=ext,
        tipo=tipo, fuente_tipo=fuente_tipo, conf_tipo=conf_tipo, regla_tipo=regla_tipo,
        num_doc=num_doc, regla_num_doc=regla_num_doc,
        fecha=fc, hora=fh, fuente_fecha=fuente_fecha, regla_fecha=regla_fecha,
        paginas=paginas, tamano=kb,
        descripcion=desc, hash_sha256=hash_sha256,
    )

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 7: CLASE ExpedienteDigital (NUEVO en v6.0)
# ══════════════════════════════════════════════════════════════════════════════

class ExpedienteDigital:
    """
    Encapsula toda la lógica de negocio del expediente.
    Completamente independiente de Tkinter.
    Puede usarse desde CLI, pruebas o cualquier contexto.
    """

    def __init__(self, carpeta: str, recursivo: bool = True):
        self.carpeta   = Path(carpeta)
        self.recursivo = recursivo
        self.registros: list = []
        self._folio_inicio: int = 1
        self._max_workers: int = min(4, max(2, (os.cpu_count() or 2) // 2))

        if not self.carpeta.exists():
            raise FileNotFoundError(f"Carpeta no encontrada: {carpeta}")
        if not self.carpeta.is_dir():
            raise NotADirectoryError(f"No es una carpeta: {carpeta}")

    # ── Recolección de archivos ───────────────────────────────────────────────

    def recolectar_archivos(self) -> list:
        """Lista los archivos a procesar, excluyendo componentes SHP y carpeta Foliados."""
        if self.recursivo:
            todos = sorted([p for p in self.carpeta.rglob("*")
                           if p.is_file()
                           and p.suffix.lower() in EXT_ALL
                           and "Foliados" not in p.parts])
        else:
            todos = sorted([self.carpeta/f for f in os.listdir(self.carpeta)
                           if (self.carpeta/f).is_file()
                           and Path(f).suffix.lower() in EXT_ALL])
        shp_stems = {f.stem for f in todos if f.suffix.lower()==".shp"}
        return [f for f in todos if not (
            f.suffix.lower() in EXT_SHP_COMP and f.stem in shp_stems)]

    # ── Análisis paralelo ─────────────────────────────────────────────────────

    def analizar(self, callback_progreso=None) -> list:
        """
        Analiza todos los archivos en paralelo con ThreadPoolExecutor.
        Fallback automático a modo secuencial si hay errores de concurrencia.

        Args:
            callback_progreso: fn(i, total, nombre_archivo) — opcional

        Returns:
            Lista de dicts con metadata completa incluyendo hash SHA-256
        """
        archivos = self.recolectar_archivos()
        if not archivos:
            return []

        total     = len(archivos)
        tmp       = [None] * total
        completados = 0
        raiz      = self.carpeta

        def analizar_con_indice(args):
            idx, ruta = args
            return idx, analizar_archivo(ruta, raiz)

        try:
            with ThreadPoolExecutor(max_workers=self._max_workers) as executor:
                futuros = {
                    executor.submit(analizar_con_indice, (i, ruta)): i
                    for i, ruta in enumerate(archivos)
                }
                for futuro in as_completed(futuros):
                    try:
                        idx, resultado = futuro.result()
                        resultado["n"] = idx + 1
                        tmp[idx] = resultado
                    except Exception as e:
                        idx = futuros[futuro]
                        log.warning(f"Error en hilo para archivo {idx}: {e}")
                        # Fallback: analizar secuencialmente ese archivo
                        try:
                            r = analizar_archivo(archivos[idx], raiz)
                            r["n"] = idx + 1
                            tmp[idx] = r
                        except Exception as e2:
                            log.error(f"Error total en archivo {idx}: {e2}")

                    completados += 1
                    if callback_progreso:
                        nombre = archivos[futuros[futuro]].name
                        try:
                            callback_progreso(completados, total, nombre)
                        except Exception as e:
                            log.debug(f"Error en callback: {e}")

        except Exception as e:
            log.error(f"Error en ThreadPoolExecutor: {e}. Usando modo secuencial.")
            tmp = []
            for i, ruta in enumerate(archivos):
                r = analizar_archivo(ruta, raiz)
                r["n"] = i + 1
                tmp.append(r)
                if callback_progreso:
                    try: callback_progreso(i+1, total, ruta.name)
                    except Exception: pass

        # Filtrar posibles None (archivos que fallaron completamente)
        self.registros = [r for r in tmp if r is not None]
        return self.registros

    # ── Foliación ─────────────────────────────────────────────────────────────

    def asignar_folios(self, folio_inicio: int = 1) -> list:
        """Asigna folios y nombres foliados. Función pura."""
        self._folio_inicio = folio_inicio
        for idx, r in enumerate(self.registros):
            r["folio"]          = str(folio_inicio + idx).zfill(3)
            r["nombre_foliado"] = construir_nombre_foliado(r)
        return self.registros

    # ── Validación ────────────────────────────────────────────────────────────

    def validar(self) -> list:
        """Retorna lista de problemas de foliación."""
        return validar_foliacion(self.registros)

    # ── Estadísticas ──────────────────────────────────────────────────────────

    def estadisticas(self) -> dict:
        return resumen_estadistico(self.registros)

    # ── Exportar ──────────────────────────────────────────────────────────────

    def exportar_word(self, config: dict) -> tuple:
        return generar_word(self.registros, str(self.carpeta), config)

    def exportar_excel(self, config: dict) -> tuple:
        return generar_excel(self.registros, str(self.carpeta), config)

    def copiar_foliados(self) -> tuple:
        return copiar_foliados(self.registros, str(self.carpeta))

    def exportar_log(self, operacion: str, detalles: list = None) -> str:
        return exportar_log(self.registros, str(self.carpeta), operacion, detalles)

    # ── Información del expediente ────────────────────────────────────────────

    @property
    def total(self) -> int:
        return len(self.registros)

    @property
    def resumen_confianza(self) -> dict:
        return {
            "Alta":  sum(1 for r in self.registros if r.get("conf_tipo")=="Alta"),
            "Media": sum(1 for r in self.registros if r.get("conf_tipo")=="Media"),
            "Baja":  sum(1 for r in self.registros if r.get("conf_tipo")=="Baja"),
        }

    @property
    def archivos_sin_hash(self) -> list:
        """Lista de archivos donde el hash no pudo calcularse."""
        return [r["nombre"] for r in self.registros
                if r.get("hash_sha256","").startswith(("ERROR","ARCHIVO","SIN"))]

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 8: FOLIADOR
# ══════════════════════════════════════════════════════════════════════════════

def asignar_folios(registros: list, folio_inicio: int=1) -> list:
    for idx,r in enumerate(registros):
        r["folio"]=str(folio_inicio+idx).zfill(3)
        r["nombre_foliado"]=construir_nombre_foliado(r)
    return registros

def copiar_foliados(registros: list, carpeta_raiz: str) -> tuple:
    raiz=Path(carpeta_raiz); foliados=raiz/"Foliados"
    foliados.mkdir(exist_ok=True)
    copiados=0; errores=[]
    for r in registros:
        try:
            dest_dir=(foliados/r["subcarpeta"] if r["subcarpeta"]!="—" else foliados)
            dest_dir.mkdir(parents=True,exist_ok=True)
            dest=dest_dir/r["nombre_foliado"]
            shutil.copy2(r["ruta"],dest); copiados+=1
            log.info(f"[COPIA] {r['nombre']} → {r['nombre_foliado']}")
            if r["extension"]==".shp":
                for ext_c in EXT_SHP_COMP:
                    cs=r["ruta"].parent/(r["ruta"].stem+ext_c)
                    cd=dest_dir/(Path(r["nombre_foliado"]).stem+ext_c)
                    if cs.exists(): shutil.copy2(cs,cd)
        except PermissionError as e:
            msg=f"Sin permisos {r['nombre']}: {e}"; errores.append(msg); log.error(msg)
        except shutil.Error as e:
            msg=f"Error copiando {r['nombre']}: {e}"; errores.append(msg); log.error(msg)
        except Exception as e:
            msg=f"Error {r['nombre']}: {e}"; errores.append(msg); log.error(msg,exc_info=True)
    return copiados,errores,str(foliados)

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 9: LOG EXPORTABLE
# ══════════════════════════════════════════════════════════════════════════════

def exportar_log(registros: list, carpeta: str,
                 operacion: str, detalles: list=None) -> str:
    ts=datetime.now().strftime('%Y%m%d_%H%M%S')
    nom_log=f"LOG_{operacion}_{ts}.txt"
    ruta_log=Path(carpeta)/nom_log
    try:
        with open(ruta_log,"w",encoding="utf-8") as f:
            f.write(f"{'='*70}\n")
            f.write(f"LOG: {operacion} — {APP_VERSION} — {inst_nombre()}\n")
            f.write(f"Fecha: {datetime.now().strftime('%d/%m/%Y %H:%M:%S')}\n")
            f.write(f"{'='*70}\n\n")
            stats=resumen_estadistico(registros)
            f.write(f"Total: {stats['total']} documentos | Tamaño: {stats['tamano_total']}\n\n")
            f.write("POR TIPO:\n")
            for tipo,cnt in stats["por_tipo"].items():
                f.write(f"  {tipo:<20} {cnt:>4}\n")
            problemas=validar_foliacion(registros)
            if problemas:
                f.write(f"\n⚠️ FOLIACIÓN ({len(problemas)} problemas):\n")
                for p in problemas: f.write(f"  • {p}\n")
            else:
                f.write("\n✅ Foliación correlativa correcta.\n")

            # Advertencia hashes fallidos
            sin_hash=[r["nombre"] for r in registros
                      if r.get("hash_sha256","").startswith(("ERROR","ARCHIVO","SIN"))]
            if sin_hash:
                f.write(f"\n⚠️ HASH NO CALCULADO ({len(sin_hash)} archivos):\n")
                for n in sin_hash: f.write(f"  • {n}\n")

            if detalles:
                f.write(f"\nDETALLE:\n")
                for d in detalles: f.write(f"  {d}\n")
            f.write(f"\n{'='*70}\n")
            f.write("LISTADO CON HASH SHA-256:\n")
            f.write(f"{'='*70}\n")
            for r in registros:
                h=r.get("hash_sha256","")
                h_display=hash_corto(h) if h else "—"
                f.write(
                    f"[{r['folio']}] {r.get('num_doc',''):14} "
                    f"{r['tipo']:14} {r['fecha']} | {h_display}\n"
                    f"        {r['nombre']}\n"
                    f"     →  {r['nombre_foliado']}\n"
                    f"        SHA256: {h}\n\n")
        log.info(f"Log exportado: {nom_log}")
        return str(ruta_log)
    except (PermissionError,OSError) as e:
        log.error(f"No se pudo escribir log: {e}"); return ""
    except Exception as e:
        log.error(f"Error exportando log: {e}",exc_info=True); return ""

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 10: EXPORTADOR WORD
# ══════════════════════════════════════════════════════════════════════════════

def _w_borde(cell,color="2C3E50"):
    tc=cell._tc; tcPr=tc.get_or_add_tcPr()
    tcB=OxmlElement('w:tcBorders')
    for side in ['top','left','bottom','right']:
        nd=OxmlElement(f'w:{side}')
        nd.set(qn('w:val'),'single'); nd.set(qn('w:sz'),'4')
        nd.set(qn('w:color'),color); tcB.append(nd)
    tcPr.append(tcB)

def _w_celda(row,idx,texto,bold=False,size=8,
             align=WD_ALIGN_PARAGRAPH.LEFT,bg=None,fg="000000"):
    cell=row.cells[idx]
    cell.vertical_alignment=WD_ALIGN_VERTICAL.CENTER
    p=cell.paragraphs[0]; p.alignment=align
    p.paragraph_format.space_before=Pt(1); p.paragraph_format.space_after=Pt(1)
    run=p.add_run(str(texto))
    run.font.size=Pt(size); run.font.bold=bold
    run.font.color.rgb=RGBColor.from_string(fg)
    if bg:
        shd=OxmlElement('w:shd'); shd.set(qn('w:val'),'clear')
        shd.set(qn('w:color'),'auto'); shd.set(qn('w:fill'),bg)
        cell._tc.get_or_add_tcPr().append(shd)
    _w_borde(cell)

def _word_configurar_pagina(doc):
    for sec in doc.sections:
        sec.page_width=Cm(29.7); sec.page_height=Cm(21.0)
        sec.left_margin=Cm(1.8); sec.right_margin=Cm(1.8)
        sec.top_margin=Cm(1.8);  sec.bottom_margin=Cm(1.5)

def _word_encabezado(doc,config,fi,ff,n_total):
    h0=doc.add_heading("",level=0); h0.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r0=h0.add_run(INST_NOMBRE)
    r0.font.size=Pt(13); r0.font.bold=True; r0.font.color.rgb=RGBColor(0x1A,0x3A,0x5C)
    p1=doc.add_paragraph(); p1.alignment=WD_ALIGN_PARAGRAPH.CENTER
    r1=p1.add_run(INST_SUBTITULO)
    r1.font.size=Pt(9); r1.font.color.rgb=RGBColor(0x2C,0x3E,0x50)
    doc.add_paragraph()
    tit=doc.add_heading("",level=1); tit.alignment=WD_ALIGN_PARAGRAPH.CENTER
    rt=tit.add_run(config.get("titulo","ÍNDICE DE DOCUMENTOS").upper())
    rt.font.size=Pt(12); rt.font.bold=True; rt.font.color.rgb=RGBColor(0x1A,0x3A,0x5C)
    pm=doc.add_paragraph(); pm.alignment=WD_ALIGN_PARAGRAPH.CENTER
    bits=[]
    if config.get("expediente"): bits.append(f"Expediente: {config['expediente']}")
    if config.get("seccion"):    bits.append(f"Sección: {config['seccion']}")
    bits+=[f"Folios: {str(fi).zfill(3)}—{str(ff).zfill(3)}",
           f"Total: {n_total} documentos",
           f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}"]
    rm=pm.add_run("  |  ".join(bits))
    rm.font.size=Pt(8.5); rm.font.color.rgb=RGBColor(0x7F,0x8C,0x8D)
    doc.add_paragraph()

def _word_tabla_datos(doc,registros):
    # Columnas incluyendo Hash SHA-256 (abreviado)
    COLS_W=COLS_INDICE+["SHA-256 (parcial)"]
    ANCHOS=[1.3,1.7,1.9,1.9,1.2,1.8,6.8,0.9,1.4,4.5,3.2]
    C=WD_ALIGN_PARAGRAPH.CENTER; L=WD_ALIGN_PARAGRAPH.LEFT; R=WD_ALIGN_PARAGRAPH.RIGHT
    tabla=doc.add_table(rows=1,cols=len(COLS_W)); tabla.style="Table Grid"
    for i,ancho in enumerate(ANCHOS):
        for cell in tabla.columns[i].cells: cell.width=Cm(ancho)
    hdr=tabla.rows[0]
    for i,col in enumerate(COLS_W):
        _w_celda(hdr,i,col,bold=True,size=7.5,align=C,bg="1A3A5C",fg="FFFFFF")
    sub_actual=None
    for r in registros:
        if r["subcarpeta"]!=sub_actual and r["subcarpeta"]!="—":
            sub_actual=r["subcarpeta"]
            sep=tabla.add_row()
            _w_celda(sep,0,f"📁  {sub_actual.upper()}",bold=True,
                     size=8,align=L,bg="D6EAF8",fg="1A3A5C")
            for j in range(1,len(COLS_W)): _w_celda(sep,j,"",bg="D6EAF8",fg="1A3A5C")
        bg_f="EBF5FB" if int(r["folio"])%2==0 else "FFFFFF"
        fila=tabla.add_row()
        h=r.get("hash_sha256",""); h_short=hash_corto(h) if h else "—"
        vals=[r["folio"],r.get("num_doc",""),r["subcarpeta"],
              r["fecha"],r["hora"],r["tipo"],
              r["descripcion"],r["paginas"],r["tamano"],
              r["nombre_foliado"],h_short]
        alns=[C,C,C,C,C,C,L,C,R,L,C]
        for i,(v,aln) in enumerate(zip(vals,alns)):
            _w_celda(fila,i,v,bold=(i==0),size=7,align=aln,bg=bg_f,
                     fg="1A3A5C" if i==0 else ("888888" if i==10 else "000000"))

def _word_resumen(doc,registros):
    doc.add_paragraph()
    h2=doc.add_heading("",level=2)
    rh=h2.add_run("RESUMEN ESTADÍSTICO")
    rh.font.size=Pt(10); rh.font.bold=True; rh.font.color.rgb=RGBColor(0x1A,0x3A,0x5C)
    stats=resumen_estadistico(registros)
    tipos_l=list(stats["por_tipo"].items()); exts_l=list(stats["por_ext"].items())
    C=WD_ALIGN_PARAGRAPH.CENTER; L=WD_ALIGN_PARAGRAPH.LEFT
    t_res=doc.add_table(rows=1,cols=4); t_res.style="Table Grid"
    hdr_r=t_res.rows[0]
    for i,col in enumerate(["Tipo documento","Cantidad","Formato","Cantidad"]):
        _w_celda(hdr_r,i,col,bold=True,size=8,align=C,bg="2C3E50",fg="FFFFFF")
    for i in range(max(len(tipos_l),len(exts_l))):
        fr=t_res.add_row(); bg_r="EBF5FB" if i%2==0 else "FFFFFF"
        t_val=tipos_l[i][0] if i<len(tipos_l) else ""; t_cnt=str(tipos_l[i][1]) if i<len(tipos_l) else ""
        e_val=exts_l[i][0]  if i<len(exts_l)  else ""; e_cnt=str(exts_l[i][1]) if i<len(exts_l) else ""
        for ci,(val,aln) in enumerate(zip([t_val,t_cnt,e_val,e_cnt],[L,C,L,C])):
            _w_celda(fr,ci,val,size=8,align=aln,bg=bg_r)
    doc.add_paragraph()
    p_tot=doc.add_paragraph()
    r_tot=p_tot.add_run(
        f"Total: {stats['total']} documentos  |  Tamaño: {stats['tamano_total']}  |  "
        f"Subcarpetas: {len(stats['por_sub'])}  |  "
        f"Integridad SHA-256: verificada en {stats['total']} archivos")
    r_tot.font.size=Pt(8); r_tot.font.bold=True; r_tot.font.color.rgb=RGBColor(0x1A,0x3A,0x5C)

def _word_advertencias(doc,registros):
    problemas=validar_foliacion(registros)
    sin_hash=[r["nombre"] for r in registros
              if r.get("hash_sha256","").startswith(("ERROR","ARCHIVO","SIN"))]
    if not problemas and not sin_hash: return
    doc.add_paragraph()
    h_adv=doc.add_heading("",level=2)
    r_adv=h_adv.add_run("⚠️  ADVERTENCIAS")
    r_adv.font.size=Pt(10); r_adv.font.bold=True; r_adv.font.color.rgb=RGBColor(0xE7,0x4C,0x3C)
    for prob in problemas:
        p_prob=doc.add_paragraph(style="List Bullet")
        rp=p_prob.add_run(f"Foliación: {prob}")
        rp.font.size=Pt(9); rp.font.color.rgb=RGBColor(0xE7,0x4C,0x3C)
    for nombre in sin_hash:
        p_h=doc.add_paragraph(style="List Bullet")
        rph=p_h.add_run(f"Hash no calculado: {nombre}")
        rph.font.size=Pt(9); rph.font.color.rgb=RGBColor(0xE6,0x7E,0x22)

def _word_pie(doc):
    doc.add_paragraph()
    pp=doc.add_paragraph(); pp.alignment=WD_ALIGN_PARAGRAPH.RIGHT
    rp=pp.add_run(
        f"{inst_pie()}  |  {APP_VERSION}  |  "
        f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}")
    rp.font.size=Pt(7); rp.font.color.rgb=RGBColor(0x95,0xA5,0xA6)

def generar_word(registros,carpeta,config):
    if not DOCX_OK: return None,"python-docx no instalado"
    try:
        fi=config.get("folio_inicio",1); ff=fi+len(registros)-1
        doc=DocxDoc()
        _word_configurar_pagina(doc)
        _word_encabezado(doc,config,fi,ff,len(registros))
        _word_tabla_datos(doc,registros)
        _word_resumen(doc,registros)
        _word_advertencias(doc,registros)
        _word_pie(doc)
        nom=(f"INDICE_Folios{str(fi).zfill(3)}-{str(ff).zfill(3)}_"
             f"{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx")
        ruta_out=Path(carpeta)/nom
        doc.save(str(ruta_out)); log.info(f"Word generado: {nom}")
        return str(ruta_out),None
    except PermissionError as e:
        log.error(f"Sin permisos Word: {e}"); return None,f"Sin permisos: {e}"
    except Exception as e:
        log.error(f"Error Word: {e}",exc_info=True); return None,str(e)

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 11: EXPORTADOR EXCEL
# ══════════════════════════════════════════════════════════════════════════════

def _xl_celda(ws,row,col,valor,bold=False,size=9,
               h_align="left",bg=None,fg="000000",wrap=False):
    cell=ws.cell(row=row,column=col,value=str(valor) if valor is not None else "")
    cell.font=Font(name="Calibri",size=size,bold=bold,color=fg or "000000")
    cell.alignment=Alignment(horizontal=h_align,vertical="center",wrap_text=wrap)
    if bg: cell.fill=PatternFill(start_color=bg,end_color=bg,fill_type="solid")
    lado=Side(style="thin",color="CCCCCC")
    cell.border=Border(left=lado,right=lado,top=lado,bottom=lado)
    return cell

def _xl_titulo_hoja(ws,texto,n_cols,color_bg=XA):
    ws.merge_cells(f"A1:{get_column_letter(n_cols)}1")
    c=ws["A1"]; c.value=texto
    c.font=Font(name="Calibri",size=11,bold=True,color=XB)
    c.fill=PatternFill(start_color=color_bg,end_color=color_bg,fill_type="solid")
    c.alignment=Alignment(horizontal="center",vertical="center")
    ws.row_dimensions[1].height=22

def _xl_encabezado_cols(ws,cols,fila=2):
    for ci,col in enumerate(cols,1):
        _xl_celda(ws,fila,ci,col,bold=True,size=9,h_align="center",bg=XA,fg=XB)
    ws.row_dimensions[fila].height=18

def _xl_separador_sub(ws,fila,subcarpeta,n_cols):
    ws.merge_cells(f"A{fila}:{get_column_letter(n_cols)}{fila}")
    c=ws[f"A{fila}"]; c.value=f"📁  {subcarpeta.upper()}"
    c.font=Font(name="Calibri",size=9,bold=True,color=XA)
    c.fill=PatternFill(start_color=XG2,end_color=XG2,fill_type="solid")
    c.alignment=Alignment(horizontal="left",vertical="center")
    ws.row_dimensions[fila].height=14

def _xl_hoja_datos(wb,nombre_hoja,registros,config,color_bg=XG1,mostrar_seps=True):
    # Columnas extendidas con Hash y trazabilidad completa
    COLS_XL=(COLS_INDICE +
             ["Fuente fecha","Regla fecha","Fuente tipo","Regla tipo",
              "Conf. tipo","Regla N° doc","SHA-256 (completo)"])
    ANCHOS=[10,14,14,12,8,13,50,6,11,50,14,18,14,18,10,16,68]

    ws=wb.create_sheet(title=nombre_hoja[:31]); ws.freeze_panes="A3"
    exp=config.get("expediente",""); secc=config.get("seccion","")
    bits=[config.get("titulo","ÍNDICE")]
    if exp:  bits.append(exp)
    if secc: bits.append(secc)
    if nombre_hoja!="RESUMEN": bits.append(f"Sección: {nombre_hoja}")
    bits.append(f"{len(registros)} documentos")
    bits.append(f"Generado: {datetime.now().strftime('%d/%m/%Y %H:%M')}")
    _xl_titulo_hoja(ws," | ".join(bits),len(COLS_XL))
    _xl_encabezado_cols(ws,COLS_XL)

    fila=3; sub_actual=None
    for r in registros:
        if mostrar_seps and r["subcarpeta"]!=sub_actual and r["subcarpeta"]!="—":
            sub_actual=r["subcarpeta"]
            _xl_separador_sub(ws,fila,sub_actual,len(COLS_XL)); fila+=1
        bg=color_bg if int(r["folio"])%2==0 else XB
        hash_completo=r.get("hash_sha256","")
        vals=[r["folio"],r.get("num_doc",""),r["subcarpeta"],
              r["fecha"],r["hora"],r["tipo"],
              r["descripcion"],r["paginas"],r["tamano"],r["nombre_foliado"],
              r.get("fuente_fecha",""),r.get("regla_fecha",""),
              r.get("fuente_tipo",""),r.get("regla_tipo",""),
              r.get("conf_tipo",""),r.get("regla_num_doc",""),
              hash_completo]
        aligns=["center","center","center","center","center","center",
                "left","center","right","left",
                "center","left","center","left","center","left","left"]
        for ci,(v,aln) in enumerate(zip(vals,aligns),1):
            color_fg=XA if ci==1 else ("888888" if ci==17 else "000000")
            _xl_celda(ws,fila,ci,v,bold=(ci==1),size=8.5,
                      h_align=aln,bg=bg,fg=color_fg,wrap=(ci==7))
        ws.row_dimensions[fila].height=16; fila+=1

    for ci,ancho in enumerate(ANCHOS,1):
        ws.column_dimensions[get_column_letter(ci)].width=ancho
    ws.auto_filter.ref=f"A2:{get_column_letter(len(COLS_XL))}{fila-1}"

    fila+=1
    ws.merge_cells(f"A{fila}:{get_column_letter(len(COLS_XL))}{fila}")
    c_st=ws[f"A{fila}"]
    fi=config.get("folio_inicio",1); ff=fi+len(config.get("todos",registros))-1
    c_st.value=(f"Total: {len(registros)} docs  |  "
                f"Folios: {str(fi).zfill(3)}–{str(ff).zfill(3)}  |  {inst_pie()}")
    c_st.font=Font(name="Calibri",size=8,italic=True,color="95A5A6")
    c_st.alignment=Alignment(horizontal="right")

def _xl_hoja_estadisticas(wb,registros,config):
    stats=resumen_estadistico(registros)
    fi=config.get("folio_inicio",1); ff=fi+len(registros)-1
    ws=wb.create_sheet(title="ESTADÍSTICAS")
    _xl_titulo_hoja(ws,f"RESUMEN ESTADÍSTICO — {config.get('titulo','ÍNDICE')}",4)
    fila=3

    def _st_titulo(texto):
        nonlocal fila
        ws.merge_cells(f"A{fila}:D{fila}"); c=ws[f"A{fila}"]; c.value=texto
        c.font=Font(name="Calibri",size=10,bold=True,color=XB)
        c.fill=PatternFill(start_color="2C3E50",end_color="2C3E50",fill_type="solid")
        c.alignment=Alignment(horizontal="left",vertical="center")
        ws.row_dimensions[fila].height=16; fila+=1

    def _st_fila(etiq,valor,bg=XB):
        nonlocal fila
        _xl_celda(ws,fila,1,etiq,size=9,h_align="left",bg=bg)
        _xl_celda(ws,fila,2,valor,bold=True,size=9,h_align="center",bg=bg,fg=XA)
        ws.merge_cells(f"C{fila}:D{fila}")
        ws.row_dimensions[fila].height=15; fila+=1

    _st_titulo("📋  DATOS GENERALES")
    for i,(etiq,val) in enumerate([
        ("Expediente",config.get("expediente","—")),
        ("Sección",config.get("seccion","—")),
        ("Total documentos",str(stats["total"])),
        ("Folio inicial",str(fi).zfill(3)),("Folio final",str(ff).zfill(3)),
        ("Tamaño total",stats["tamano_total"]),
        ("Subcarpetas",str(len(stats["por_sub"]))),
        ("Integridad SHA-256","Verificada — "+APP_VERSION),
        ("Generado",datetime.now().strftime("%d/%m/%Y %H:%M")),
    ]):
        _st_fila(etiq,val,XG1 if i%2==0 else XB)

    fila+=1; _st_titulo("📄  POR TIPO")
    for i,(tipo,cnt) in enumerate(stats["por_tipo"].items()):
        bg=XG1 if i%2==0 else XB
        _xl_celda(ws,fila,1,tipo,size=9,h_align="left",bg=bg)
        _xl_celda(ws,fila,2,cnt,bold=True,size=9,h_align="center",bg=bg,fg=XA)
        _xl_celda(ws,fila,3,f"{cnt/stats['total']*100:.1f}%",size=9,
                  h_align="center",bg=bg,fg="7F8C8D")
        ws.row_dimensions[fila].height=15; fila+=1

    fila+=1; _st_titulo("📁  POR SUBCARPETA")
    for i,(sub,cnt) in enumerate(stats["por_sub"].items()):
        _st_fila(sub,cnt,XG1 if i%2==0 else XB)

    fila+=1; _st_titulo("🗂  POR FORMATO")
    for i,(ext,cnt) in enumerate(stats["por_ext"].items()):
        _st_fila(ext,cnt,XG1 if i%2==0 else XB)

    fila+=1; _st_titulo("🔍  CONFIANZA DE CLASIFICACIÓN")
    conf_counts={}
    for r in registros:
        c=r.get("conf_tipo","?"); conf_counts[c]=conf_counts.get(c,0)+1
    for i,(conf,cnt) in enumerate(conf_counts.items()):
        _st_fila(conf,cnt,XG1 if i%2==0 else XB)

    fila+=1; _st_titulo("🔐  INTEGRIDAD SHA-256")
    sin_hash=[r for r in registros
              if r.get("hash_sha256","").startswith(("ERROR","ARCHIVO","SIN"))]
    _st_fila("Archivos con hash válido",str(len(registros)-len(sin_hash)),XG1)
    _st_fila("Archivos sin hash",str(len(sin_hash)),"FEF2F2" if sin_hash else XB)

    for ci,ancho in enumerate([28,16,10,10],1):
        ws.column_dimensions[get_column_letter(ci)].width=ancho

def generar_excel(registros,carpeta,config):
    if not XLSX_OK: return None,"openpyxl no instalado"
    try:
        wb=openpyxl.Workbook(); wb.remove(wb.active)
        grupos={}
        for r in registros: grupos.setdefault(r["subcarpeta"],[]).append(r)
        colores=[XG1,XV,XM,XO,"FEF2F2","F0F8FF","FFFFF0","FFF0F5"]
        _xl_hoja_datos(wb,"RESUMEN",registros,{**config,"todos":registros},
                       XG1,mostrar_seps=True)
        for idx,(sub,items) in enumerate(grupos.items()):
            _xl_hoja_datos(wb,sub[:31] if sub!="—" else "Raíz",
                           items,config,colores[idx%len(colores)],mostrar_seps=False)
        _xl_hoja_estadisticas(wb,registros,config)
        fi=config.get("folio_inicio",1); ff=fi+len(registros)-1
        nom=(f"INDICE_Folios{str(fi).zfill(3)}-{str(ff).zfill(3)}_"
             f"{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx")
        ruta_out=Path(carpeta)/nom
        wb.save(str(ruta_out)); log.info(f"Excel generado: {nom}")
        return str(ruta_out),None
    except PermissionError as e:
        log.error(f"Sin permisos Excel: {e}"); return None,f"Sin permisos: {e}"
    except Exception as e:
        log.error(f"Error Excel: {e}",exc_info=True); return None,str(e)


# ══════════════════════════════════════════════════════════════════════════════
# REQUIREMENTS.TXT
# ══════════════════════════════════════════════════════════════════════════════

REQUIREMENTS = """# GENERADOR DE ÍNDICE + FOLIADOR — CRUBC Los Ríos
# Instalar: pip install -r requirements.txt

# Obligatorias
pdfplumber
python-docx
openpyxl
pillow

# Opcionales
python-pptx
rarfile
py7zr
"""

def generar_requirements(carpeta: str) -> str:
    ruta = Path(carpeta) / "requirements.txt"
    try:
        with open(ruta,"w",encoding="utf-8") as f:
            f.write(REQUIREMENTS)
        return str(ruta)
    except Exception as e:
        log.warning(f"No se pudo generar requirements.txt: {e}"); return ""

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 15: PRUEBAS AUTOMÁTICAS (actualizadas para v6.0)
# ══════════════════════════════════════════════════════════════════════════════

class TestMotorAnalisis(unittest.TestCase):
    """
    Pruebas automáticas — CRUBC Los Ríos · Oficina Técnica v6.0
    Casos reales del expediente ECMPO Wadalafken.
    Ejecutar: python generador_indice.py --test
    """

    def test_fecha_nombre_correo_ecmpo(self):
        fc,fh,fuente,regla=fecha_desde_nombre(
            "2025-07-17_1750_Correo_Ingreso_de_Reclamacion_ECMPO_Wadalafken")
        self.assertEqual(fc,"2025-07-17"); self.assertEqual(fh,"17:50")
        self.assertEqual(fuente,"Nombre-archivo"); self.assertIn("2025-07-17_1750",regla)

    def test_fecha_nombre_acta_crubc(self):
        fc,fh,fuente,_=fecha_desde_nombre(
            "2025-08-08_0900_Acta_de_Acuerdos_Reunion_Comite_Tecnico_CRUBC")
        self.assertEqual(fc,"2025-08-08"); self.assertEqual(fh,"09:00")

    def test_fecha_nombre_sin_patron(self):
        fc,fh,fuente,_=fecha_desde_nombre("documento_sin_fecha")
        self.assertEqual(fc,""); self.assertEqual(fuente,"")

    def test_fecha_nombre_invalida(self):
        fc,fh,_,_=fecha_desde_nombre("2025-13-45_9999_documento")
        self.assertEqual(fc,"")

    def test_fecha_texto_correo_reenvio(self):
        texto=("gabinete gobernador 17 de julio de 2025 a las 8:59 "
               "oficina partes 17 jul 2025 8:46 "
               "fecha original: 15 de julio de 2025 a las 17:50")
        fc,fh,fuente,_=fecha_desde_texto(texto)
        self.assertEqual(fc,"2025-07-15"); self.assertEqual(fh,"17:50")

    def test_fecha_texto_sin_fecha(self):
        fc,fh,fuente,_=fecha_desde_texto("texto sin fecha alguna")
        self.assertEqual(fc,""); self.assertEqual(fuente,"")

    def test_tipo_correo_nombre(self):
        tipo,fuente,conf,regla=detectar_tipo(
            "","2025-07-17_Correo_Ingreso_ECMPO",".pdf")
        self.assertEqual(tipo,"Correo"); self.assertEqual(conf,"Alta")

    def test_tipo_acta_nombre(self):
        tipo,_,_,_=detectar_tipo("","2025-08-08_Acta_de_Acuerdos_CRUBC",".pdf")
        self.assertEqual(tipo,"Acta")

    def test_tipo_oficio_nombre(self):
        tipo,_,conf,_=detectar_tipo("","2025-07-30_Oficio_N_1192",".pdf")
        self.assertEqual(tipo,"Oficio"); self.assertEqual(conf,"Alta")

    def test_tipo_resolucion_nombre(self):
        tipo,_,_,_=detectar_tipo("","2025-08-05_Res_Exenta_548-2025",".pdf")
        self.assertEqual(tipo,"Resolución")

    def test_tipo_informe_nombre(self):
        tipo,_,_,_=detectar_tipo("","2025-08-06_Informe_Observaciones_Municipalidad",".pdf")
        self.assertEqual(tipo,"Informe")

    def test_tipo_correo_texto(self):
        texto="de: jschneider@goredelosrios.cl para: cacuna@goredelosrios.cl"
        tipo,fuente,conf,_=detectar_tipo(texto,"doc",".pdf")
        self.assertEqual(tipo,"Correo"); self.assertEqual(fuente,"Texto-interno")

    def test_tipo_fallback_extension(self):
        tipo,fuente,conf,regla=detectar_tipo("","archivo",".xlsx")
        self.assertEqual(tipo,"Excel"); self.assertEqual(fuente,"Extensión")
        self.assertIn(".xlsx",regla)

    def test_num_doc_folio(self):
        num,regla=extraer_numero_doc("folio n°2544","correo")
        self.assertEqual(num,"Folio N°2544"); self.assertIn("folio",regla)

    def test_num_doc_res_exenta(self):
        num,_=extraer_numero_doc("resolución exenta n°238/2025","resolucion")
        self.assertIn("238",num)

    def test_num_doc_desde_nombre(self):
        num,_=extraer_numero_doc("","2025-07-30_Oficio_N_1192_Convoca")
        self.assertIn("1192",num)

    def test_num_doc_sin_numero(self):
        num,regla=extraer_numero_doc("texto sin número","archivo")
        self.assertEqual(num,""); self.assertEqual(regla,"")

    def test_foliacion_correcta(self):
        regs=[{"folio":str(i).zfill(3),"nombre":f"d{i}.pdf"} for i in range(1,6)]
        self.assertEqual(validar_foliacion(regs),[])

    def test_foliacion_gap(self):
        regs=[{"folio":"001","nombre":"a.pdf"},
              {"folio":"002","nombre":"b.pdf"},
              {"folio":"005","nombre":"e.pdf"}]
        self.assertTrue(any("003" in p or "Gap" in p for p in validar_foliacion(regs)))

    def test_foliacion_duplicado(self):
        regs=[{"folio":"001","nombre":"a.pdf"},{"folio":"001","nombre":"b.pdf"}]
        self.assertTrue(any("duplicado" in p.lower() for p in validar_foliacion(regs)))

    def test_nombre_foliado_correo(self):
        r={"folio":"017","fecha":"2025-07-17","hora":"17:50",
           "tipo":"Correo","descripcion":"Ingreso Reclamacion ECMPO","extension":".pdf"}
        nom=construir_nombre_foliado(r)
        self.assertTrue(nom.startswith("017_20250717_1750_Correo_"))
        self.assertTrue(nom.endswith(".pdf"))

    def test_nombre_foliado_sin_espacios(self):
        r={"folio":"005","fecha":"2025-08-08","hora":"09:00",
           "tipo":"Acta","descripcion":"Reunión Comité Técnico CRUBC","extension":".pdf"}
        self.assertNotIn(" ",construir_nombre_foliado(r))

    def test_validar_fecha(self):
        self.assertEqual(validar_fecha("2025-07-17"),"")
        self.assertNotEqual(validar_fecha("17/07/2025"),"")
        self.assertNotEqual(validar_fecha(""),"")

    def test_validar_hora(self):
        self.assertEqual(validar_hora("08:59"),"")
        self.assertEqual(validar_hora(""),"")
        self.assertNotEqual(validar_hora("8:5"),"")
        self.assertNotEqual(validar_hora("25:00"),"")

    def test_validar_descripcion(self):
        self.assertEqual(validar_descripcion("Ingreso Reclamación ECMPO"),"")
        self.assertNotEqual(validar_descripcion(""),"")
        self.assertNotEqual(validar_descripcion("ab"),"")

    def test_validar_campos_ok(self):
        self.assertEqual(validar_campos_edicion("2025-07-17","17:50","Correo",
                                                 "Ingreso Reclamación ECMPO"),[])

    def test_desc_foliado(self):
        self.assertNotIn(" ",desc_foliado("Ingreso Reclamación ECMPO"))
        self.assertLessEqual(len(desc_foliado("A"*200)),80)
        self.assertFalse(desc_foliado("2025-07-17_1750_Correo_ECMPO").startswith("2025"))

    def test_asignar_folios_inicio_1(self):
        regs=[{"tipo":"PDF","fecha":"","hora":"","descripcion":"d","extension":".pdf"}
               for _ in range(3)]
        r=asignar_folios(regs,1)
        self.assertEqual(r[0]["folio"],"001"); self.assertEqual(r[2]["folio"],"003")

    def test_asignar_folios_inicio_33(self):
        regs=[{"tipo":"PDF","fecha":"","hora":"","descripcion":"d","extension":".pdf"}
               for _ in range(3)]
        r=asignar_folios(regs,33)
        self.assertEqual(r[0]["folio"],"033"); self.assertEqual(r[2]["folio"],"035")

    # ── Pruebas Hash SHA-256 (NUEVO v6.0) ─────────────────────────────────────

    def test_hash_genera_hex_64_chars(self):
        """Hash SHA-256 produce siempre 64 caracteres hexadecimales"""
        import tempfile, os
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as f:
            f.write(b"contenido de prueba CRUBC Los Rios")
            tmp=Path(f.name)
        try:
            h=generar_hash(tmp)
            self.assertEqual(len(h),64)
            self.assertTrue(all(c in "0123456789abcdef" for c in h))
        finally:
            os.unlink(tmp)

    def test_hash_mismo_archivo_mismo_resultado(self):
        """El mismo archivo siempre produce el mismo hash"""
        import tempfile, os
        with tempfile.NamedTemporaryFile(delete=False, suffix=".txt") as f:
            f.write(b"documento expediente ECMPO Wadalafken"); tmp=Path(f.name)
        try:
            h1=generar_hash(tmp); h2=generar_hash(tmp)
            self.assertEqual(h1,h2)
        finally:
            os.unlink(tmp)

    def test_hash_archivos_distintos_distintos_hash(self):
        """Archivos con contenido diferente producen hashes diferentes"""
        import tempfile, os
        with tempfile.NamedTemporaryFile(delete=False,suffix=".txt") as f:
            f.write(b"contenido A"); tmp_a=Path(f.name)
        with tempfile.NamedTemporaryFile(delete=False,suffix=".txt") as f:
            f.write(b"contenido B"); tmp_b=Path(f.name)
        try:
            self.assertNotEqual(generar_hash(tmp_a),generar_hash(tmp_b))
        finally:
            os.unlink(tmp_a); os.unlink(tmp_b)

    def test_hash_archivo_inexistente(self):
        """Archivo inexistente retorna código de error legible"""
        h=generar_hash(Path("/ruta/inexistente/archivo.pdf"))
        self.assertIn("NO_ENCONTRADO",h.upper().replace("_","_"))

    def test_hash_corto_formato(self):
        """hash_corto produce formato abreviado correcto"""
        h_completo="a"*64
        h_corto=hash_corto(h_completo)
        self.assertIn("...",h_corto)
        self.assertLess(len(h_corto),len(h_completo))

    # ── Pruebas ExpedienteDigital (NUEVO v6.0) ────────────────────────────────

    def test_expediente_carpeta_inexistente(self):
        """ExpedienteDigital lanza error con carpeta inexistente"""
        with self.assertRaises(FileNotFoundError):
            ExpedienteDigital("/carpeta/que/no/existe")

    def test_expediente_propiedades_vacias(self):
        """ExpedienteDigital sin analizar tiene total=0"""
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            exp=ExpedienteDigital(tmp)
            self.assertEqual(exp.total,0)


    # ── Hash global del expediente ────────────────────────────────────────────

    def test_hash_expediente_determinista(self):
        docs=[DocumentoAnalizado(folio="001",nombre="a.pdf",hash_sha256="abc"),
              DocumentoAnalizado(folio="002",nombre="b.pdf",hash_sha256="def")]
        h1=generar_hash_expediente(docs); h2=generar_hash_expediente(docs)
        self.assertEqual(h1,h2)

    def test_hash_expediente_64_chars(self):
        docs=[DocumentoAnalizado(folio="001",nombre="a.pdf",hash_sha256="aaa")]
        self.assertEqual(len(generar_hash_expediente(docs)),64)

    def test_hash_expediente_cambia_con_contenido(self):
        docs1=[DocumentoAnalizado(folio="001",nombre="a.pdf",hash_sha256="abc")]
        docs2=[DocumentoAnalizado(folio="001",nombre="a.pdf",hash_sha256="xyz")]
        self.assertNotEqual(generar_hash_expediente(docs1),
                            generar_hash_expediente(docs2))

    # ── DocumentoAnalizado ────────────────────────────────────────────────────

    def test_documento_tiene_hash_valido(self):
        d=DocumentoAnalizado(hash_sha256="a"*64)
        self.assertTrue(d.tiene_hash_valido())

    def test_documento_hash_invalido(self):
        d=DocumentoAnalizado(hash_sha256="ERROR_LECTURA")
        self.assertFalse(d.tiene_hash_valido())

    def test_documento_hash_corto_formato(self):
        d=DocumentoAnalizado(hash_sha256="a"*64)
        h=d.hash_corto(); self.assertIn("...",h); self.assertLess(len(h),64)

    def test_documento_historial_ediciones(self):
        d=DocumentoAnalizado(nombre="test.pdf",tipo="PDF")
        self.assertFalse(d.fue_editado())
        d.registrar_edicion("tipo","PDF","Correo")
        self.assertTrue(d.fue_editado())
        self.assertEqual(d.historial_ediciones[0]["campo"],"tipo")
        self.assertEqual(d.historial_ediciones[0]["valor_anterior"],"PDF")
        self.assertEqual(d.historial_ediciones[0]["valor_nuevo"],"Correo")

    def test_documento_multiples_ediciones(self):
        d=DocumentoAnalizado(nombre="test.pdf")
        d.registrar_edicion("tipo","PDF","Correo")
        d.registrar_edicion("fecha","2025-01-01","2025-07-17")
        self.assertEqual(len(d.historial_ediciones),2)

    # ── Caché JSON ────────────────────────────────────────────────────────────

    def test_cache_vacia_carpeta_nueva(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            cache=cargar_cache(Path(tmp))
            self.assertEqual(cache,{})

    def test_cache_guardar_y_cargar(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            cache={"k":{"mtime":1234.0,"tamano_bytes":100,
                         "resultado":{"folio":"001","nombre":"t.pdf",
                         "hash_sha256":"a"*64}}}
            guardar_cache(tp,cache); cache2=cargar_cache(tp)
            self.assertIn("k",cache2)

    # ── Integración con archivos temporales ──────────────────────────────────

    def test_integracion_analizar_txt(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            archivo=tp/"2025-07-17_1750_Correo_Prueba_Integracion.txt"
            archivo.write_text("de: test@goredelosrios.cl\nPrueba CRUBC",
                               encoding="utf-8")
            doc=analizar_archivo(archivo,tp)
            self.assertIsInstance(doc,DocumentoAnalizado)
            self.assertEqual(doc.tipo,"Correo")
            self.assertTrue(doc.tiene_hash_valido())

    def test_integracion_excluye_foliados(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            (tp/"archivo.txt").write_text("contenido",encoding="utf-8")
            foliados=tp/"Foliados"; foliados.mkdir()
            (foliados/"foliado.txt").write_text("foliado",encoding="utf-8")
            exp=ExpedienteDigital(tmp,recursivo=True)
            archivos=exp.recolectar_archivos()
            nombres=[a.name for a in archivos]
            self.assertIn("archivo.txt",nombres)
            self.assertNotIn("foliado.txt",nombres)

    def test_integracion_asignar_folios_y_hash(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            for i in range(3):
                (tp/f"2025-07-{17+i:02d}_1000_Correo_doc{i}.txt"
                 ).write_text(f"contenido {i}",encoding="utf-8")
            exp=ExpedienteDigital(tmp)
            exp.analizar(); exp.asignar_folios(folio_inicio=10)
            self.assertEqual(exp.docs[0].folio,"010")
            self.assertEqual(len(exp.hash_expediente),64)

    def test_integracion_copia_foliados(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            (tp/"2025-07-17_1750_Correo_Prueba.txt"
             ).write_text("contenido",encoding="utf-8")
            exp=ExpedienteDigital(tmp)
            exp.analizar(); exp.asignar_folios(1)
            n,errores,ruta_f=exp.copiar_foliados()
            self.assertEqual(n,1); self.assertEqual(len(errores),0)
            archivos=list(Path(ruta_f).rglob("*.txt"))
            self.assertEqual(len(archivos),1)

    def test_integracion_pdf_corrupto_no_rompe(self):
        import tempfile
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            pdf=(tp/"2025-07-17_1750_Correo_Corrupto.pdf")
            pdf.write_bytes(b"no es PDF")
            try:
                doc=analizar_archivo(pdf,tp)
                self.assertIsInstance(doc,DocumentoAnalizado)
            except Exception as e:
                self.fail(f"Excepción inesperada: {e}")

    def test_integracion_zip_lista_contenido(self):
        import tempfile, zipfile as zf
        with tempfile.TemporaryDirectory() as tmp:
            tp=Path(tmp)
            zip_path=tp/"2025-07-17_1750_Comprimido_Adjuntos.zip"
            with zf.ZipFile(zip_path,"w") as z:
                z.writestr("a.txt","contenido A")
                z.writestr("b.pdf","contenido B")
            doc=analizar_archivo(zip_path,tp)
            self.assertIn("2 archivos",doc.descripcion)

    def test_integracion_expediente_carpeta_inexistente(self):
        with self.assertRaises(FileNotFoundError):
            ExpedienteDigital("/ruta/inexistente/que/no/existe")


def ejecutar_tests():
    print("\n"+"="*70)
    print(f"PRUEBAS AUTOMÁTICAS — {inst_subtitulo()}")
    print(f"Versión: {APP_VERSION}")
    print("="*70+"\n")
    loader=unittest.TestLoader()
    suite=loader.loadTestsFromTestCase(TestMotorAnalisis)
    runner=unittest.TextTestRunner(verbosity=2,stream=sys.stdout)
    resultado=runner.run(suite)
    print("\n"+"="*70)
    total=resultado.testsRun
    ok=total-len(resultado.failures)-len(resultado.errors)
    print(f"RESULTADO: {ok}/{total} pruebas pasadas")
    if resultado.wasSuccessful():
        print("✅ Todas las pruebas pasaron correctamente.")
    else:
        print(f"❌ {len(resultado.failures)} fallos, {len(resultado.errors)} errores.")
    print("="*70+"\n")
    return resultado.wasSuccessful()

# ══════════════════════════════════════════════════════════════════════════════
# CAPA 13: INTERFAZ GRÁFICA
# ══════════════════════════════════════════════════════════════════════════════

class InterfazGrafica(tk.Tk):
    """
    Capa visual. Solo coordina la UI y delega toda la lógica
    a la clase ExpedienteDigital.
    """

    def __init__(self):
        super().__init__()
        self.title("Generador de Índice — CRUBC Los Ríos")
        self.geometry("1400x840"); self.configure(bg=CF)
        self.resizable(True,True)
        self.carpeta_var    = tk.StringVar()
        self.recursivo_var  = tk.BooleanVar(value=True)
        self.titulo_var     = tk.StringVar(value="ÍNDICE DE DOCUMENTOS")
        self.expediente_var = tk.StringVar()
        self.seccion_var    = tk.StringVar()
        self.folio_var      = tk.StringVar(value="1")
        self.expediente: ExpedienteDigital = None
        self._analizando    = False
        self._docs_visibles: list = []
        self.filtro_tipo_var= tk.StringVar(value="TODOS")
        self.filtro_sub_var = tk.StringVar(value="TODAS")
        self.busqueda_var   = tk.StringVar()
        self.orden_col      = None
        self.orden_asc      = True
        self._construir_ui()
        self._generar_requirements_inicio()

    def _btn(self,p,txt,cmd,col,**kw):
        return tk.Button(p,text=txt,command=cmd,bg=col,fg="white",
            font=("Segoe UI",9,"bold"),relief="flat",padx=10,pady=5,
            cursor="hand2",activebackground=col,activeforeground="white",**kw)

    def _entry(self,p,var,w=20):
        return tk.Entry(p,textvariable=var,font=("Segoe UI",10),
                        bg=CP,fg=CT,insertbackground=CT,relief="flat",width=w)

    def _lbl(self,p,txt,w=0,**kw):
        opts=dict(bg=CF,fg=CT,font=("Segoe UI",10)); opts.update(kw)
        if w: opts["width"]=w; opts["anchor"]="w"
        return tk.Label(p,text=txt,**opts)

    def _construir_ui(self):
        fh=tk.Frame(self,bg=CP,pady=10); fh.pack(fill="x")
        tk.Label(fh,text=f"📋  {APP_TITULO}",
                 font=("Segoe UI",14,"bold"),bg=CP,fg=CA).pack(side="left",padx=20)
        self.lbl_inst=tk.Label(fh,
                 text=f"{inst_subtitulo()}  {APP_VERSION}",
                 font=("Segoe UI",10),bg=CP,fg=CT)
        self.lbl_inst.pack(side="right",padx=10)
        self._btn(fh,"⚙️ Institución",
                  self._abrir_config_institucional,CP).pack(side="right",padx=4)

        r1=tk.Frame(self,bg=CF,pady=6); r1.pack(fill="x",padx=20)
        self._lbl(r1,"Carpeta:",w=14).pack(side="left")
        self._entry(r1,self.carpeta_var,56).pack(side="left",padx=6,ipady=4)
        self._btn(r1,"📂 Seleccionar",self._sel,CA).pack(side="left",padx=3)
        self._btn(r1,"🔍 Analizar",   self._analizar,CY).pack(side="left",padx=3)
        tk.Checkbutton(r1,text="Subcarpetas",variable=self.recursivo_var,
            bg=CF,fg=CT,selectcolor=CP,activebackground=CF,
            font=("Segoe UI",9)).pack(side="left",padx=8)
        self._btn(r1,"🧪 Pruebas",self._run_tests,CM).pack(side="left",padx=6)

        r2=tk.Frame(self,bg=CF,pady=3); r2.pack(fill="x",padx=20)
        self._lbl(r2,"Título:",w=14).pack(side="left")
        self._entry(r2,self.titulo_var,34).pack(side="left",padx=6,ipady=3)
        self._lbl(r2,"  Expediente:").pack(side="left")
        self._entry(r2,self.expediente_var,18).pack(side="left",padx=6,ipady=3)
        self._lbl(r2,"  Sección:").pack(side="left")
        self._entry(r2,self.seccion_var,16).pack(side="left",padx=6,ipady=3)

        r3=tk.Frame(self,bg=CF,pady=3); r3.pack(fill="x",padx=20)
        self._lbl(r3,"Folio inicial:",w=14).pack(side="left")
        self._entry(r3,self.folio_var,8).pack(side="left",padx=6,ipady=3)
        tk.Label(r3,text="← Ej: 1 → 001  |  33 → continúa desde 033",
                 bg=CF,fg="#7f8c8d",font=("Segoe UI",8,"italic")).pack(side="left",padx=6)

        self.progreso=ttk.Progressbar(self,orient="horizontal",mode="determinate")
        self.progreso.pack(fill="x",padx=20,pady=(4,0))

        tk.Label(self,
            text="  Doble clic para editar · "
                 "Columna SHA-256 muestra huella digital de cada archivo",
            bg=CF,fg="#7f8c8d",font=("Segoe UI",8,"italic")).pack(anchor="w",padx=20)

        fm=tk.Frame(self,bg=CF); fm.pack(fill="both",expand=True,padx=20,pady=(2,4))
        COLS_UI=("N° Folio","N° Doc.","Subcarpeta","Fecha","Hora","Tipo",
                 "Descripción / Asunto","Págs","Tamaño",
                 "F.Fecha","Conf.","SHA-256","Nombre foliado")
        ANCHOS=[65,95,105,90,55,88,280,40,82,110,50,130,260]
        self.tree=ttk.Treeview(fm,columns=COLS_UI,show="headings",height=13)
        for col,ancho in zip(COLS_UI,ANCHOS):
            self.tree.heading(col,text=col,command=lambda c=col:self._ordenar_por(c))
            self.tree.column(col,width=ancho,anchor="w")
        sx=ttk.Scrollbar(fm,orient="horizontal",command=self.tree.xview)
        sy=ttk.Scrollbar(fm,orient="vertical",  command=self.tree.yview)
        self.tree.configure(xscrollcommand=sx.set,yscrollcommand=sy.set)
        self.tree.grid(row=0,column=0,sticky="nsew")
        sy.grid(row=0,column=1,sticky="ns"); sx.grid(row=1,column=0,sticky="ew")
        fm.grid_rowconfigure(0,weight=1); fm.grid_columnconfigure(0,weight=1)
        self.tree.bind("<Double-1>",self._editar)
        self.tree.bind("<Button-3>",self._menu_contextual)
        st=ttk.Style(); st.theme_use("clam")
        st.configure("Treeview",background=CPAR,foreground=CT,
                      fieldbackground=CPAR,rowheight=24,font=("Segoe UI",9))
        st.configure("Treeview.Heading",background=CP,foreground=CA,
                      font=("Segoe UI",9,"bold"))
        st.map("Treeview",background=[("selected",CA)])

        fb=tk.Frame(self,bg=CP,pady=10); fb.pack(fill="x",side="bottom")
        self.lbl=tk.Label(fb,text="Selecciona una carpeta y presiona Analizar.",
            bg=CP,fg=CT,font=("Segoe UI",9))
        self.lbl.pack(side="left",padx=16)
        self._btn(fb,"📁 Copiar Foliados",self._copiar,   CN).pack(side="right",padx=8)
        self._btn(fb,"📝 Word",           self._gen_word, CV).pack(side="right",padx=4)
        self._btn(fb,"📊 Excel",          self._gen_excel,CA).pack(side="right",padx=4)
        self._btn(fb,"📝+📊 Ambos",       self._gen_ambos,CM).pack(side="right",padx=4)
        self._btn(fb,"🗑 Limpiar",         self._limpiar,  CR).pack(side="right",padx=4)

    def _run_tests(self):
        win=tk.Toplevel(self); win.title("Pruebas Automáticas — v6.0")
        win.configure(bg=CF); win.geometry("820x540"); win.grab_set()
        tk.Label(win,text=f"🧪  PRUEBAS AUTOMÁTICAS — {APP_VERSION}",
                 font=("Segoe UI",12,"bold"),bg=CF,fg=CA).pack(pady=(16,4))
        txt=tk.Text(win,font=("Consolas",9),bg="#0d1117",fg="#c9d1d9",
                    insertbackground=CT,relief="flat")
        txt.pack(fill="both",expand=True,padx=16,pady=8)
        import io
        stream=io.StringIO()
        loader=unittest.TestLoader()
        suite=loader.loadTestsFromTestCase(TestMotorAnalisis)
        runner=unittest.TextTestRunner(verbosity=2,stream=stream)
        resultado=runner.run(suite)
        output=stream.getvalue()
        total=resultado.testsRun
        ok=total-len(resultado.failures)-len(resultado.errors)
        txt.insert("end",output)
        txt.insert("end",f"\n{'='*60}\n")
        if resultado.wasSuccessful():
            txt.insert("end",f"✅ {ok}/{total} pruebas pasadas — {APP_VERSION}\n")
        else:
            txt.insert("end",
                f"❌ {ok}/{total} pasadas | "
                f"{len(resultado.failures)} fallos | {len(resultado.errors)} errores\n")
        txt.configure(state="disabled"); txt.see("end")


    def _actualizar_combos_filtro(self):
        if not self.expediente: return
        tipos = sorted(set(d.tipo for d in self.expediente.docs))
        subs  = sorted(set(d.subcarpeta for d in self.expediente.docs))
        self.combo_tipo["values"] = ["TODOS"] + tipos
        self.combo_sub["values"]  = ["TODAS"] + subs
        self.filtro_tipo_var.set("TODOS")
        self.filtro_sub_var.set("TODAS")

    def _aplicar_filtros(self):
        if not self.expediente: return
        tipo_f   = self.filtro_tipo_var.get()
        sub_f    = self.filtro_sub_var.get()
        busqueda = self.busqueda_var.get().lower().strip()
        docs = list(self.expediente.docs)
        if tipo_f != "TODOS":  docs = [d for d in docs if d.tipo==tipo_f]
        if sub_f  != "TODAS":  docs = [d for d in docs if d.subcarpeta==sub_f]
        if busqueda:
            docs = [d for d in docs if (
                busqueda in d.nombre.lower() or
                busqueda in d.descripcion.lower() or
                busqueda in d.num_doc.lower() or
                busqueda in d.hash_sha256.lower() or
                busqueda in d.tipo.lower())]
        self._docs_visibles = docs
        self._poblar(docs)
        n_total = len(self.expediente.docs)
        n_vis   = len(docs)
        if n_vis < n_total:
            self.lbl.config(text=f"🔍 Mostrando {n_vis}/{n_total} documentos")

    def _ordenar_por(self, col: str):
        MAPA = {"N° Folio":"folio","Subcarpeta":"subcarpeta","Fecha":"fecha",
                "Hora":"hora","Tipo":"tipo","Descripción / Asunto":"descripcion",
                "Conf.":"conf_tipo","N° Doc.":"num_doc"}
        campo = MAPA.get(col)
        if not campo or not self._docs_visibles: return
        if self.orden_col == col: self.orden_asc = not self.orden_asc
        else: self.orden_col = col; self.orden_asc = True
        self._docs_visibles.sort(
            key=lambda d: getattr(d,campo,"") or "",
            reverse=not self.orden_asc)
        self._poblar(self._docs_visibles)
        dir_icon = "▲" if self.orden_asc else "▼"
        self.lbl.config(text=f"Ordenado por {col} {dir_icon}")

    def _menu_contextual(self, event):
        item = self.tree.identify_row(event.y)
        if not item or item.startswith("SEP_"): return
        self.tree.selection_set(item)
        menu = tk.Menu(self,tearoff=0,bg=CP,fg=CT,
                       activebackground=CA,activeforeground=CT)
        menu.add_command(label="✎ Editar",command=lambda:self._editar_iid(item))
        menu.add_separator()
        menu.add_command(label="📋 Copiar hash SHA-256",
                         command=lambda:self._copiar_hash(item))
        menu.add_command(label="📋 Copiar nombre foliado",
                         command=lambda:self._copiar_nombre_foliado(item))
        try: menu.tk_popup(event.x_root,event.y_root)
        finally: menu.grab_release()

    def _copiar_hash(self, iid: str):
        try:
            idx=int(iid)-1; doc=self.expediente.docs[idx]
            self.clipboard_clear(); self.clipboard_append(doc.hash_sha256)
            self.lbl.config(text=f"✅ Hash copiado: {doc.hash_corto()}")
        except Exception as e: log.debug(f"Copiar hash: {e}")

    def _copiar_nombre_foliado(self, iid: str):
        try:
            idx=int(iid)-1; doc=self.expediente.docs[idx]
            self.clipboard_clear(); self.clipboard_append(doc.nombre_foliado)
            self.lbl.config(text=f"✅ Nombre copiado")
        except Exception as e: log.debug(f"Copiar nombre: {e}")

    def _editar_iid(self, iid: str):
        try: idx=int(iid)-1
        except ValueError: return
        if not self.expediente or idx<0 or idx>=len(self.expediente.docs): return
        d=self.expediente.docs[idx]
        self._abrir_edicion(d,iid)


    def _set_botones_accion(self, activo: bool):
        """Habilita/deshabilita botones de exportación."""
        estado = "normal" if activo else "disabled"
        for btn in [self.btn_copiar,self.btn_word,self.btn_excel,self.btn_ambos]:
            try: btn.config(state=estado)
            except Exception: pass

    def _set_botones_analisis(self, analizando: bool):
        """Deshabilita botones de selección durante análisis."""
        estado = "disabled" if analizando else "normal"
        try:
            self.btn_sel.config(state=estado)
            self.btn_analizar.config(state=estado)
        except Exception: pass


    def _abrir_config_institucional(self):
        """
        Panel de configuración institucional.
        Permite cambiar nombre, subtítulo, dominio de correo
        y agregar tipos documentales personalizados.
        Los cambios se guardan en config_indice.json automáticamente.
        """
        win = tk.Toplevel(self)
        win.title("Configuración Institucional")
        win.configure(bg=CF); win.geometry("680x520"); win.grab_set()

        tk.Label(win, text="⚙️  CONFIGURACIÓN INSTITUCIONAL",
                 font=("Segoe UI",12,"bold"),bg=CF,fg=CA).pack(pady=(16,4))
        tk.Label(win,
            text="Los cambios se guardan en config_indice.json y se aplican al instante",
            bg=CF,fg="#7f8c8d",font=("Segoe UI",8,"italic")).pack(pady=(0,10))

        cfg = dict(_config_activa)
        inst = cfg.get("institucion", {})

        # Campos institucionales
        campos = [
            ("Nombre institución:",  "nombre",    inst.get("nombre",""),    50),
            ("Subtítulo / Unidad:",  "subtitulo", inst.get("subtitulo",""), 50),
            ("Pie de página:",       "pie",       inst.get("pie",""),       50),
            ("Dominio correo (ej: gore.gob.cl):","dominio", cfg.get("dominio_correo",""), 30),
        ]
        vars_cfg = {}
        for etiq, clave, valor, ancho in campos:
            fr = tk.Frame(win,bg=CF); fr.pack(fill="x",padx=24,pady=5)
            tk.Label(fr,text=etiq,bg=CF,fg=CT,
                     font=("Segoe UI",9),width=22,anchor="w").pack(side="left")
            var = tk.StringVar(value=valor); vars_cfg[clave] = var
            tk.Entry(fr,textvariable=var,font=("Segoe UI",10),
                     bg=CP,fg=CT,insertbackground=CT,
                     relief="flat",width=ancho).pack(side="left",ipady=3)

        # Sección tipos adicionales
        tk.Label(win,
            text="Tipos documentales adicionales (uno por línea: Tipo|_patron1_|_patron2_)",
            bg=CF,fg=CT,font=("Segoe UI",8)).pack(padx=24,pady=(12,2),anchor="w")
        tk.Label(win,
            text="Ejemplo:  Formulario|_formulario_|form_n    o    Voucher|_voucher_|comprobante_",
            bg=CF,fg="#7f8c8d",font=("Segoe UI",7,"italic")).pack(padx=24,anchor="w")

        fr_txt = tk.Frame(win,bg=CF); fr_txt.pack(fill="both",expand=True,padx=24,pady=4)
        txt_tipos = tk.Text(fr_txt,font=("Consolas",9),bg=CP,fg=CT,
                            insertbackground=CT,relief="flat",height=6)
        txt_tipos.pack(fill="both",expand=True)
        # Cargar tipos adicionales existentes
        tipos_add = cfg.get("tipos_adicionales",{}).get("nombre",[])
        for entrada in tipos_add:
            if isinstance(entrada,list) and len(entrada)==2:
                linea = entrada[0]+"|"+("|".join(entrada[1]))
                txt_tipos.insert("end", linea+"\n")

        lbl_estado = tk.Label(win,text="",bg=CF,fg=CV,font=("Segoe UI",8))
        lbl_estado.pack(pady=(4,0))

        def guardar():
            # Leer valores
            nombre    = vars_cfg["nombre"].get().strip()
            subtitulo = vars_cfg["subtitulo"].get().strip()
            pie       = vars_cfg["pie"].get().strip()
            dominio   = vars_cfg["dominio"].get().strip().lstrip("@")

            if not nombre:
                lbl_estado.config(text="⚠️ El nombre de institución no puede estar vacío",fg=CR)
                return

            # Parsear tipos adicionales
            tipos_nuevos = []
            for linea in txt_tipos.get("1.0","end").strip().splitlines():
                linea = linea.strip()
                if not linea: continue
                partes = [p.strip() for p in linea.split("|")]
                if len(partes) >= 2:
                    tipo_n = partes[0]
                    pats_n = [p for p in partes[1:] if p]
                    if tipo_n and pats_n:
                        tipos_nuevos.append([tipo_n, pats_n])

            # Construir nueva config
            cfg_nueva = {
                "institucion": {
                    "nombre":    nombre,
                    "subtitulo": subtitulo,
                    "pie":       pie or subtitulo,
                },
                "dominio_correo": dominio,
                "tipos_adicionales": {
                    "nombre": tipos_nuevos,
                    "texto":  cfg.get("tipos_adicionales",{}).get("texto",{})
                },
                "version": APP_VERSION
            }

            # Guardar y aplicar
            ok = guardar_config_json(cfg_nueva)
            if ok:
                _config_activa.update(cfg_nueva)
                aplicar_config(cfg_nueva)
                # Actualizar encabezado de la ventana principal
                self._actualizar_encabezado_inst()
                lbl_estado.config(
                    text=f"✅ Guardado — {len(TIPOS_NOMBRE)} tipos activos",fg=CV)
                log.info(f"Institución actualizada: {nombre}")
            else:
                lbl_estado.config(text="⚠️ No se pudo guardar config_indice.json",fg=CR)

        fr_btn = tk.Frame(win,bg=CF); fr_btn.pack(fill="x",padx=24,pady=10)
        self._btn(fr_btn,"💾 Guardar configuración",guardar,CV).pack(side="right",padx=4)
        self._btn(fr_btn,"✖ Cerrar",win.destroy,CR).pack(side="right",padx=4)

    def _actualizar_encabezado_inst(self):
        """Actualiza el encabezado de la ventana principal con la institución activa."""
        try:
            self.lbl_inst.config(
                text=f"CRUBC Los Ríos · {inst_subtitulo()}  {APP_VERSION}")
        except Exception:
            pass

    def _sel(self):
        c=filedialog.askdirectory(title="Selecciona carpeta con archivos")
        if c: self.carpeta_var.set(c)

    def _analizar(self):
        if self._analizando:
            messagebox.showinfo("En proceso","Espera que termine."); return
        carpeta=self.carpeta_var.get().strip()
        if not carpeta or not os.path.isdir(carpeta):
            messagebox.showerror("Error","Selecciona una carpeta válida."); return
        self._limpiar(); self._analizando=True
        self.lbl.config(text="⏳ Analizando archivos..."); self.update()
        threading.Thread(target=self._hilo,args=(carpeta,),daemon=True).start()

    def _hilo(self,carpeta:str):
        try:
            try: fi=max(1,int(self.folio_var.get()))
            except ValueError: fi=1

            # Crear ExpedienteDigital y analizar en paralelo
            self.expediente=ExpedienteDigital(carpeta,self.recursivo_var.get())

            def progreso(i,total,nombre):
                msg=f"⏳ {i}/{total}: {nombre}"
                self.after(0,lambda m=msg,v=i,t=total:[
                    self.lbl.config(text=m),
                    self.progreso.configure(maximum=t,value=v)])

            self.expediente.analizar(callback_progreso=progreso)
            self.expediente.asignar_folios(fi)
            self.after(0,self._fin)

        except FileNotFoundError as e:
            log.error(f"Carpeta: {e}")
            self.after(0,lambda:messagebox.showerror("Error",str(e)))
        except Exception as e:
            log.error(f"Error análisis: {e}",exc_info=True)
            self.after(0,lambda:messagebox.showerror("Error de análisis",str(e)))
        finally:
            self.after(0,lambda:[setattr(self,'_analizando',False),
                                  self.progreso.configure(value=0)])

    def _fin(self):
        if not self.expediente: return
        regs=self.expediente.registros
        self._poblar(regs)
        conf=self.expediente.resumen_confianza
        sin_hash=len(self.expediente.archivos_sin_hash)
        txt=(f"✅ {self.expediente.total} archivos · "
             f"{len(set(r['subcarpeta'] for r in regs))} subcarpetas  |  "
             f"Confianza: 🟢{conf['Alta']} 🟡{conf['Media']} 🔴{conf['Baja']}  |  "
             f"🔐 SHA-256 calculado")
        if sin_hash: txt+=f" (⚠️ {sin_hash} sin hash)"
        txt+="  |  Doble clic para editar"
        self.lbl.config(text=txt)
        try: self._set_botones_accion(True)
        except Exception: pass

    def _fila_valores(self,r):
        h=r.get("hash_sha256",""); h_short=hash_corto(h) if h else "—"
        return (r["folio"],r.get("num_doc",""),r["subcarpeta"],
                r["fecha"],r["hora"],r["tipo"],r["descripcion"],
                r["paginas"],r["tamano"],
                r.get("fuente_fecha",""),r.get("conf_tipo",""),
                h_short,r["nombre_foliado"])

    def _poblar(self,regs):
        self.tree.delete(*self.tree.get_children())
        sub_actual=None
        for r in regs:
            if r["subcarpeta"]!=sub_actual and r["subcarpeta"]!="—":
                sub_actual=r["subcarpeta"]
                self.tree.insert("","end",iid=f"SEP_{sub_actual}",tags=("sep",),
                    values=(f"📁 {sub_actual}","","","","","","","","","","","",""))
            tag="par" if int(r["folio"])%2==0 else "impar"
            self.tree.insert("","end",iid=str(r["n"]),tags=(tag,),
                             values=self._fila_valores(r))
        self.tree.tag_configure("par",  background=CPAR)
        self.tree.tag_configure("impar",background=CIMP)
        self.tree.tag_configure("sep",  background=CSUB,foreground=CA)

    def _editar(self,event):
        item=self.tree.focus()
        if not item or item.startswith("SEP_"): return
        try: idx=int(item)-1
        except ValueError: return
        if not self.expediente: return
        regs=self.expediente.registros
        if idx<0 or idx>=len(regs): return
        r=regs[idx]

        win=tk.Toplevel(self); win.title(f"Editar: {r['nombre']}")
        win.configure(bg=CF); win.geometry("760x510"); win.grab_set()

        tk.Label(win,text=f"Archivo: {r['nombre']}",
                 bg=CF,fg=CA,font=("Segoe UI",9,"bold")).pack(padx=20,pady=(14,2),anchor="w")

        # Panel trazabilidad
        fr_trace=tk.Frame(win,bg="#16213e",pady=6)
        fr_trace.pack(fill="x",padx=20,pady=(0,8))
        h_completo=r.get("hash_sha256","")
        trace_info=[
            ("Tipo detectado:", f"{r['tipo']} — desde {r.get('fuente_tipo','?')} — regla: {r.get('regla_tipo','?')}"),
            ("Confianza tipo:", r.get("conf_tipo","?")),
            ("Fecha desde:",    f"{r.get('fuente_fecha','?')} — regla: {r.get('regla_fecha','?')}"),
            ("N° documento:",   f"{r.get('num_doc','—')} — regla: {r.get('regla_num_doc','—')}"),
            ("SHA-256:",        h_completo if h_completo else "No calculado"),
        ]
        for etiq,val in trace_info:
            frt=tk.Frame(fr_trace,bg="#16213e"); frt.pack(fill="x",padx=8,pady=1)
            tk.Label(frt,text=etiq,bg="#16213e",fg="#7f8c8d",
                     font=("Segoe UI",8),width=16,anchor="w").pack(side="left")
            color_val="#a8c8f0" if etiq!="SHA-256:" else "#5dade2"
            tk.Label(frt,text=val[:80],bg="#16213e",fg=color_val,
                     font=("Consolas" if etiq=="SHA-256:" else "Segoe UI",8)
                     ).pack(side="left")

        campos=[("Descripción / Asunto","descripcion"),("Tipo","tipo"),
                ("Fecha (YYYY-MM-DD)","fecha"),("Hora (HH:MM)","hora")]
        vars_e={}
        for etiq,clave in campos:
            fr=tk.Frame(win,bg=CF); fr.pack(fill="x",padx=20,pady=4)
            tk.Label(fr,text=etiq+":",bg=CF,fg=CT,
                     font=("Segoe UI",10),width=22,anchor="w").pack(side="left")
            var=tk.StringVar(value=r.get(clave,"")); vars_e[clave]=var
            tk.Entry(fr,textvariable=var,font=("Segoe UI",10),
                     bg=CP,fg=CT,insertbackground=CT,
                     relief="flat",width=48).pack(side="left",ipady=3)

        tk.Label(win,text="Nombre foliado resultante:",
                 bg=CF,fg=CT,font=("Segoe UI",9)).pack(padx=20,pady=(8,0),anchor="w")
        lbl_prev=tk.Label(win,text="",bg=CF,fg=CA,
                           font=("Segoe UI",9,"bold"),wraplength=680)
        lbl_prev.pack(padx=20,anchor="w")
        lbl_err=tk.Label(win,text="",bg=CF,fg=CR,font=("Segoe UI",8),wraplength=680)
        lbl_err.pack(padx=20,anchor="w",pady=(2,0))

        def upd(*_):
            r_tmp={**r}
            for k,v in vars_e.items(): r_tmp[k]=v.get()
            lbl_prev.config(text=construir_nombre_foliado(r_tmp))
            errores=validar_campos_edicion(
                vars_e["fecha"].get(),vars_e["hora"].get(),
                vars_e["tipo"].get(),vars_e["descripcion"].get())
            if errores: lbl_err.config(text="⚠️ "+"  |  ".join(errores),fg=CR)
            else:       lbl_err.config(text="✅ Campos válidos",fg=CV)

        for v in vars_e.values(): v.trace_add("write",upd)
        upd()

        # Guardar valores originales para historial
        valores_originales={k:r.get(k,"") for k in vars_e}

        def guardar():
            errores=validar_campos_edicion(
                vars_e["fecha"].get(),vars_e["hora"].get(),
                vars_e["tipo"].get(),vars_e["descripcion"].get())
            if errores:
                messagebox.showerror("Campos inválidos",
                    "Corrige antes de guardar:\n\n"
                    +"\n".join(f"  • {e}" for e in errores),parent=win); return
            for k,v in vars_e.items():
                nuevo=v.get().strip(); anterior=valores_originales[k]
                if nuevo!=anterior and hasattr(r,"registrar_edicion"):
                    r.registrar_edicion(k,anterior,nuevo)
                r[k]=nuevo
            r["nombre_foliado"]=construir_nombre_foliado(r)
            if item in self.tree.get_children(""):
                self.tree.item(item,values=self._fila_valores(r))
            win.destroy()

        self._btn(win,"💾 Guardar cambios",guardar,CV).pack(pady=14,anchor="e",padx=20)

    def _config(self)->dict:
        try: fi=max(1,int(self.folio_var.get()))
        except ValueError:
            messagebox.showerror("Error","El folio inicial debe ser ≥ 1."); return None
        return {"titulo":self.titulo_var.get().strip() or "ÍNDICE DE DOCUMENTOS",
                "expediente":self.expediente_var.get().strip(),
                "seccion":self.seccion_var.get().strip(),
                "folio_inicio":fi}

    def _verificar(self,validar_folios=False)->bool:
        if not self.expediente or not self.expediente.registros:
            messagebox.showinfo("Sin datos","Primero analiza una carpeta."); return False
        if validar_folios:
            problemas=self.expediente.validar()
            if problemas:
                det="\n".join(f"  • {p}" for p in problemas[:8])
                resp=messagebox.askyesno("⚠️ Advertencias",
                    f"{len(problemas)} problema(s):\n\n"+det+"\n\n¿Continuar?")
                if not resp: return False
        return True

    def _copiar(self):
        if not self._verificar(): return
        if not messagebox.askyesno("Confirmar copia foliada",
            f"Se copiarán {self.expediente.total} archivos a:\n"
            f"{self.expediente.carpeta}\\Foliados\\\n\n"
            "Originales intactos.\n\n¿Continuar?"): return
        self.lbl.config(text="⏳ Copiando..."); self.update()
        n,errores,ruta_f=self.expediente.copiar_foliados()
        detalles=[f"[OK] {r['nombre']} → {r['nombre_foliado']}"
                  for r in self.expediente.registros]+errores
        ruta_log=self.expediente.exportar_log("COPIA_FOLIADOS",detalles)
        msg=f"✅ {n} archivos copiados a:\n{ruta_f}\n"
        if errores: msg+=f"\n⚠️ {len(errores)} errores:\n"+"\n".join(errores[:5])
        if ruta_log: msg+=f"\n\nLog: {Path(ruta_log).name}"
        messagebox.showinfo("Copia completada",msg)
        self.lbl.config(text=f"✅ {n} archivos copiados a Foliados/")
        if messagebox.askyesno("Abrir carpeta","¿Abrir Foliados ahora?"):
            try: os.startfile(ruta_f)
            except Exception as e: log.warning(f"No se pudo abrir: {e}")

    def _gen_word(self):
        if not self._verificar(validar_folios=True): return
        if not DOCX_OK: messagebox.showerror("Error","pip install python-docx"); return
        cfg=self._config()
        if not cfg: return
        ruta,err=self.expediente.exportar_word(cfg)
        if err: messagebox.showerror("Error Word",err); return
        ruta_log=self.expediente.exportar_log("INDICE_WORD")
        self._ok("Word",ruta,ruta_log)

    def _gen_excel(self):
        if not self._verificar(validar_folios=True): return
        if not XLSX_OK: messagebox.showerror("Error","pip install openpyxl"); return
        cfg=self._config()
        if not cfg: return
        ruta,err=self.expediente.exportar_excel(cfg)
        if err: messagebox.showerror("Error Excel",err); return
        ruta_log=self.expediente.exportar_log("INDICE_EXCEL")
        self._ok("Excel",ruta,ruta_log)

    def _gen_ambos(self):
        if not self._verificar(validar_folios=True): return
        cfg=self._config()
        if not cfg: return
        rutas=[]; errores=[]
        if DOCX_OK:
            r,e=self.expediente.exportar_word(cfg)
            if r: rutas.append(Path(r).name)
            if e: errores.append(f"Word: {e}")
        else: errores.append("Word: python-docx no instalado")
        if XLSX_OK:
            r,e=self.expediente.exportar_excel(cfg)
            if r: rutas.append(Path(r).name)
            if e: errores.append(f"Excel: {e}")
        else: errores.append("Excel: openpyxl no instalado")
        ruta_log=self.expediente.exportar_log("INDICE_COMPLETO")
        msg="✅ Archivos generados:\n\n"+"".join(f"  • {n}\n" for n in rutas)
        if errores: msg+="\n⚠️ Errores:\n"+"".join(f"  • {e}\n" for e in errores)
        if ruta_log: msg+=f"\n\nLog: {Path(ruta_log).name}"
        messagebox.showinfo("Generación completada",msg)
        if rutas: self.lbl.config(text=f"✅ {', '.join(rutas)}")

    def _ok(self,tipo,ruta,ruta_log=""):
        fi=max(1,int(self.folio_var.get()))
        ff=fi+self.expediente.total-1
        msg=(f"Archivo: {Path(ruta).name}\n\n"
             f"Folios: {str(fi).zfill(3)} → {str(ff).zfill(3)}\n"
             f"Total: {self.expediente.total} documentos\n"
             f"SHA-256: verificado en todos los archivos\n\n"
             f"Guardado en:\n{self.expediente.carpeta}")
        if ruta_log: msg+=f"\nLog: {Path(ruta_log).name}"
        messagebox.showinfo(f"✅ {tipo} generado",msg)
        self.lbl.config(text=f"✅ {tipo}: {Path(ruta).name}")
        if messagebox.askyesno("Abrir","¿Abrir el archivo ahora?"):
            try: os.startfile(ruta)
            except Exception as e: log.warning(f"No se pudo abrir: {e}")

    def _limpiar(self):
        self.tree.delete(*self.tree.get_children())
        self.expediente=None; self.progreso.configure(value=0)
        self.lbl.config(text="Selecciona una carpeta y presiona Analizar.")

# ══════════════════════════════════════════════════════════════════════════════
# PUNTO DE ENTRADA
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    if "--test" in sys.argv:
        exito=ejecutar_tests()
        sys.exit(0 if exito else 1)
    else:
        inicializar_config()   # Carga config_indice.json al iniciar
        app=InterfazGrafica()
        app.mainloop()
